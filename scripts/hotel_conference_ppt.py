#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""高端奢华酒店会议服务手册 PPT 生成脚本"""

import os
from PIL import Image, ImageDraw, ImageFont, ImageFilter
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import math

OUTPUT_DIR = os.path.join(os.path.dirname(os.path.abspath(__file__)), "..", "outputs")
PPTX_DIR = os.path.join(OUTPUT_DIR, "pptx")
IMG_DIR = os.path.join(OUTPUT_DIR, "images")
os.makedirs(PPTX_DIR, exist_ok=True)
os.makedirs(IMG_DIR, exist_ok=True)

# ── 色彩方案 ──
GOLD = (200, 169, 110)
DARK_GOLD = (168, 140, 85)
LIGHT_GOLD = (232, 213, 163)
NAVY = (27, 27, 47)
DARK = (44, 44, 44)
LIGHT_BG = (250, 248, 245)
CREAM = (245, 240, 232)
WHITE = (255, 255, 255)
WARM_GRAY = (180, 175, 165)

# pptx 颜色对象
PPTX_GOLD = RGBColor(0xC8, 0xA9, 0x6E)
PPTX_NAVY = RGBColor(0x1B, 0x1B, 0x2F)
PPTX_DARK = RGBColor(0x2C, 0x2C, 0x2C)
PPTX_CREAM = RGBColor(0xFA, 0xF8, 0xF5)
PPTX_WHITE = RGBColor(0xFF, 0xFF, 0xFF)
PPTX_LIGHT_GOLD = RGBColor(0xE8, 0xD5, 0xA3)
PPTX_WARM_GRAY = RGBColor(0xB4, 0xAF, 0xA5)

# ── 图片生成 ──

def find_font(size=24, bold=False):
    """查找系统可用的中文字体"""
    candidates = [
        "C:/Windows/Fonts/msyh.ttc",
        "C:/Windows/Fonts/msyhbd.ttc",
        "C:/Windows/Fonts/simhei.ttf",
        "C:/Windows/Fonts/simsun.ttc",
        "C:/Windows/Fonts/STXINWEI.TTF",
        "C:/Windows/Fonts/FZNEWSHU.TTF",
    ]
    if bold:
        candidates = [
            "C:/Windows/Fonts/msyhbd.ttc",
            "C:/Windows/Fonts/simhei.ttf",
            "C:/Windows/Fonts/msyh.ttc",
        ] + candidates
    for p in candidates:
        if os.path.exists(p):
            try:
                return ImageFont.truetype(p, size)
            except:
                continue
    return ImageFont.load_default()


def draw_gold_line(draw, x, y, width, height=3):
    """绘制金色装饰线"""
    draw.rectangle([x, y, x + width, y + height], fill=GOLD)


def draw_golden_pattern(draw, x, y, width, height):
    """绘制金色装饰花纹 (简化菱形链条)"""
    step = 20
    for i in range(0, width, step):
        cx = x + i
        cy = y + height // 2
        r = 4
        draw.ellipse([cx - r, cy - r, cx + r, cy + r], outline=GOLD, width=1)


def gen_cover_image():
    """封面主图 — 金色线条勾勒的酒店会议空间"""
    w, h = 1920, 1080
    img = Image.new("RGB", (w, h), NAVY)
    draw = ImageDraw.Draw(img)

    # 抽象几何背景
    for i in range(0, w + h, 40):
        alpha = max(0, 30 - i // 30)
        if i % 80 == 0:
            draw.line([(i, 0), (0, i)], fill=(40, 40, 65), width=1)
        else:
            draw.line([(i, 0), (w, i - w)], fill=(35, 35, 55), width=1)

    # 中央建筑感抽象图形
    cx, cy = w // 2, h // 2 - 40
    # 大弧顶
    for r in range(200, 260, 10):
        draw.arc([cx - r, cy - r, cx + r, cy + r], -30, 210, fill=GOLD + (max(0, 80 - r) if r > 230 else 150,))

    # 金色竖线阵列
    for i, x_offset in enumerate(range(-180, 200, 30)):
        alpha = 180 - abs(i - 6) * 20
        if alpha > 30:
            draw.line([(cx + x_offset, cy + 60), (cx + x_offset, cy + 260)], fill=GOLD + (alpha,), width=2)

    # 金色横线
    draw_gold_line(draw, cx - 250, cy + 270, 500, 2)

    # 左下角金色装饰
    draw_golden_pattern(draw, 60, h - 80, 300, 20)

    # 右下装饰
    draw.rectangle([w - 260, 30, w - 30, 33], fill=GOLD)
    draw.rectangle([w - 260, 36, w - 180, 39], fill=GOLD)

    path = os.path.join(IMG_DIR, "cover_main.png")
    img.save(path, quality=95)
    return path


def gen_banquet_hall():
    """大宴会厅 — 金色灯光璀璨效果"""
    w, h = 1280, 800
    img = Image.new("RGB", (w, h), (20, 20, 38))
    draw = ImageDraw.Draw(img)

    # 天花板水晶灯效果
    cx, cy = w // 2, 120
    for r in range(40, 220, 15):
        color = (
            min(255, 180 + r // 2),
            min(255, 160 + r // 3),
            min(255, 100 + r // 4),
        )
        draw.ellipse([cx - r, cy - r * 0.6, cx + r, cy + r * 0.6], outline=color, width=2)

    # 金色光晕
    for r2 in range(260, 340, 20):
        alpha = max(0, 80 - (r2 - 260) // 2)
        draw.ellipse([cx - r2, cy - r2 * 0.4, cx + r2, cy + r2 * 0.4], outline=GOLD + (alpha,), width=1)

    # 柱子
    for px in [180, w - 180]:
        draw.rectangle([px - 15, 200, px + 15, h - 60], fill=(40, 38, 50))
        draw.rectangle([px - 18, 200, px + 18, 203], fill=GOLD)
        draw.rectangle([px - 18, h - 63, px + 18, h - 60], fill=GOLD)

    # 地面反射
    for i in range(0, w, 50):
        draw.line([(i, h - 100), (i + 100, h)], fill=(35, 33, 48), width=1)

    # 讲台
    draw.rectangle([cx - 80, h - 180, cx + 80, h - 60], fill=(25, 25, 42))
    draw.rectangle([cx - 85, h - 185, cx + 85, h - 182], fill=GOLD)

    # 金色装饰文字
    font_en = find_font(20)
    draw.text((cx - 80, h - 140), "GRAND BALLROOM", fill=GOLD, font=font_en)

    # 标题线
    draw_gold_line(draw, 60, h - 40, 200, 2)

    path = os.path.join(IMG_DIR, "banquet_hall.png")
    img.save(path, quality=95)
    return path


def gen_meeting_room():
    """中型会议室 — 现代商务风格"""
    w, h = 1280, 800
    img = Image.new("RGB", (w, h), (30, 30, 48))
    draw = ImageDraw.Draw(img)

    # 会议桌（椭圆）
    cx, cy = w // 2, h // 2
    draw.ellipse([cx - 300, cy - 120, cx + 300, cy + 120], fill=(20, 20, 35), outline=GOLD, width=2)

    # 椅子（小圆）
    chair_positions = []
    for angle in range(0, 360, 30):
        rad = math.radians(angle)
        rx = 350 * 0.85  # x radius
        ry = 170 * 0.85  # y radius
        px = cx + rx * math.cos(rad)
        py = cy + ry * math.sin(rad)
        draw.ellipse([px - 18, py - 18, px + 18, py + 18], fill=(50, 48, 65), outline=GOLD, width=1)
        draw.ellipse([px - 10, py - 10, px + 10, py + 10], fill=(60, 58, 75))

    # 投影幕布
    draw.rectangle([cx - 200, 60, cx + 200, 200], fill=(15, 15, 28))
    draw.rectangle([cx - 205, 55, cx + 205, 58], fill=GOLD)
    draw.rectangle([cx - 205, 197, cx + 205, 200], fill=GOLD)
    # 幕布上的"画面"
    for yy in range(70, 190, 15):
        draw.rectangle([cx - 180, yy, cx - 180 + 120, yy + 8], fill=(40, 38, 58))

    # 侧墙灯带
    draw.rectangle([30, 100, 35, h - 100], fill=GOLD + (100,))
    draw.rectangle([w - 35, 100, w - 30, h - 100], fill=GOLD + (100,))

    path = os.path.join(IMG_DIR, "meeting_room.png")
    img.save(path, quality=95)
    return path


def gen_tea_break():
    """茶歇服务 — 精致摆台"""
    w, h = 1280, 800
    img = Image.new("RGB", (w, h), (50, 44, 38))
    draw = ImageDraw.Draw(img)

    # 长桌
    draw.rectangle([150, h - 280, w - 150, h - 100], fill=(70, 62, 52))
    draw.rectangle([148, h - 285, w - 148, h - 282], fill=GOLD)
    draw.rectangle([148, h - 102, w - 148, h - 99], fill=GOLD)

    # 桌布褶皱
    for x in range(180, w - 150, 40):
        draw.line([(x, h - 280), (x + 10, h - 100)], fill=(60, 54, 44), width=1)

    # 杯子
    cup_positions = [280, 400, 520, 640, 760, 880, 1000]
    for cx in cup_positions:
        # 杯身
        draw.rectangle([cx - 15, h - 240, cx + 15, h - 200], fill=(200, 195, 185))
        draw.rectangle([cx - 17, h - 242, cx + 17, h - 239], fill=GOLD)
        # 咖啡/茶
        draw.rectangle([cx - 12, h - 235, cx + 12, h - 215], fill=(80, 50, 35))
        # 碟子
        draw.ellipse([cx - 28, h - 198, cx + 28, h - 182], fill=(210, 205, 195))
        draw.ellipse([cx - 30, h - 200, cx + 30, h - 197], fill=GOLD)

    # 装饰花
    for fx in [220, 1060]:
        for y_off in range(0, 80, 20):
            draw.ellipse([fx - 5, h - 200 - y_off, fx + 5, h - 190 - y_off], fill=(220, 180, 140))

    path = os.path.join(IMG_DIR, "tea_break.png")
    img.save(path, quality=95)
    return path


def gen_business_center():
    """商务中心"""
    w, h = 1280, 800
    img = Image.new("RGB", (w, h), (28, 30, 42))
    draw = ImageDraw.Draw(img)

    # 工作台
    draw.rectangle([100, h - 200, w - 100, h - 60], fill=(38, 40, 55))

    # 电脑屏幕
    for i, (sx, sy) in enumerate([(250, h - 380), (550, h - 380), (850, h - 380)]):
        draw.rectangle([sx, sy, sx + 160, sy + 120], fill=(15, 18, 30))
        draw.rectangle([sx - 3, sy - 3, sx + 163, sy + 123], outline=GOLD + (80,), width=1)
        # 屏幕光
        draw.rectangle([sx + 10, sy + 10, sx + 150, sy + 110], fill=(25, 30, 50))
        draw.rectangle([sx + 50, sy + 50, sx + 110, sy + 60], fill=(60, 60, 80))
        # 底座
        draw.rectangle([sx + 60, sy + 125, sx + 100, sy + 135], fill=(50, 50, 65))

    # 打印机示意
    draw.rectangle([w - 200, h - 220, w - 120, h - 160], fill=(45, 45, 60))
    draw.rectangle([w - 205, h - 225, w - 115, h - 222], fill=GOLD)

    # 装饰线
    draw_gold_line(draw, 60, 50, 300, 2)

    path = os.path.join(IMG_DIR, "business_center.png")
    img.save(path, quality=95)
    return path


def gen_concierge():
    """专属管家服务"""
    w, h = 1280, 800
    img = Image.new("RGB", (w, h), (26, 26, 40))
    draw = ImageDraw.Draw(img)

    # 人物剪影
    # 管家（居中）
    cx = w // 2
    # 身体
    draw.rectangle([cx - 40, h - 320, cx + 40, h - 100], fill=(40, 38, 55))
    # 头部
    draw.ellipse([cx - 30, h - 370, cx + 30, h - 310], fill=(45, 43, 60))
    # 领结
    draw.polygon([(cx - 8, h - 290), (cx, h - 280), (cx + 8, h - 290)], fill=GOLD)

    # 两侧人物
    for side, offset in [(-1, 180), (1, 180)]:
        px = cx + side * offset
        draw.rectangle([px - 30, h - 280, px + 30, h - 100], fill=(35, 33, 50))
        draw.ellipse([px - 22, h - 330, px + 22, h - 285], fill=(40, 38, 55))

    # 金色光芒线条
    for angle in range(0, 360, 15):
        rad = math.radians(angle)
        r1, r2 = 200, 240
        x1 = cx + r1 * math.cos(rad)
        y1 = h // 2 - 100 + r1 * math.sin(rad) * 0.5
        x2 = cx + r2 * math.cos(rad)
        y2 = h // 2 - 100 + r2 * math.sin(rad) * 0.5
        draw.line([(x1, y1), (x2, y2)], fill=GOLD + (60,), width=1)

    # 五星
    star_centers = [(cx - 120, h - 350), (cx + 120, h - 350), (cx, h - 390)]
    for sx, sy in star_centers:
        for r in range(3, 7, 1):
            draw.ellipse([sx - r, sy - r, sx + r, sy + r], outline=GOLD, width=1)

    path = os.path.join(IMG_DIR, "concierge.png")
    img.save(path, quality=95)
    return path


def gen_success_case():
    """成功案例"""
    w, h = 1280, 800
    img = Image.new("RGB", (w, h), (20, 20, 35))
    draw = ImageDraw.Draw(img)

    # 奖杯抽象造型
    cx = w // 2
    # 底座
    draw.rectangle([cx - 60, h - 120, cx + 60, h - 80], fill=GOLD)
    draw.rectangle([cx - 40, h - 160, cx + 40, h - 120], fill=GOLD + (180,))
    # 柱身
    draw.rectangle([cx - 15, h - 280, cx + 15, h - 160], fill=GOLD)
    # 杯身
    draw.ellipse([cx - 60, h - 380, cx + 60, h - 270], fill=GOLD + (200,))
    draw.ellipse([cx - 50, h - 370, cx + 50, h - 280], fill=(30, 28, 48))
    # 把手
    draw.arc([cx - 80, h - 360, cx - 55, h - 290], 270, 450, fill=GOLD, width=4)
    draw.arc([cx + 55, h - 360, cx + 80, h - 290], 90, 270, fill=GOLD, width=4)

    # 星星装饰
    for sx, sy in [(cx - 180, h - 300), (cx + 180, h - 300), (cx - 150, h - 200), (cx + 150, h - 200)]:
        for r in range(2, 5):
            draw.ellipse([sx - r, sy - r, sx + r, sy + r], outline=GOLD + (100,), width=1)

    draw_gold_line(draw, cx - 200, 60, 400, 2)

    path = os.path.join(IMG_DIR, "success_case.png")
    img.save(path, quality=95)
    return path


def gen_vip_lounge():
    """VIP休息室"""
    w, h = 1280, 800
    img = Image.new("RGB", (w, h), (38, 32, 28))
    draw = ImageDraw.Draw(img)

    # 沙发（抽象）
    draw.rectangle([200, h - 300, 500, h - 100], fill=(60, 52, 44))
    draw.rectangle([195, h - 305, 505, h - 302], fill=GOLD)
    # 靠背
    draw.rectangle([200, h - 380, 500, h - 300], fill=(55, 48, 40))
    # 坐垫线
    draw.rectangle([210, h - 280, 490, h - 275], fill=GOLD + (80,))

    # 另一侧沙发
    draw.rectangle([780, h - 300, 1080, h - 100], fill=(60, 52, 44))
    draw.rectangle([775, h - 305, 1085, h - 302], fill=GOLD)
    draw.rectangle([780, h - 380, 1080, h - 300], fill=(55, 48, 40))

    # 茶几
    draw.ellipse([cx := w // 2 - 80, h - 240, cx + 160, h - 160], fill=(45, 38, 32))
    draw.ellipse([cx - 3, h - 243, cx + 163, h - 157], outline=GOLD, width=1)

    # 落地灯
    draw.line([(w - 150, h - 100), (w - 150, h - 450)], fill=GOLD + (120,), width=3)
    draw.ellipse([w - 170, h - 480, w - 130, h - 440], fill=GOLD + (100,))
    for rr in range(60, 120, 10):
        draw.ellipse([w - 150 - rr, h - 520, w - 150 + rr, h - 440], outline=GOLD + (20,), width=1)

    path = os.path.join(IMG_DIR, "vip_lounge.png")
    img.save(path, quality=95)
    return path


# ── PPT 生成 ──

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

SLIDE_W = prs.slide_width
SLIDE_H = prs.slide_height


def add_bg(slide, color=PPTX_CREAM):
    """设置幻灯片背景"""
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_shape_bg(slide, color=PPTX_NAVY):
    """全幅深色背景矩形"""
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, SLIDE_H)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def add_gold_bar(slide, left, top, width, height=Pt(3)):
    """金色装饰条"""
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = PPTX_GOLD
    shape.line.fill.background()
    # 移到最前需要多次操作
    # 用 sp_tree 重新排序
    return shape


def add_text_box(slide, left, top, width, height, text, font_size=18,
                 bold=False, color=PPTX_DARK, align=PP_ALIGN.LEFT,
                 font_name="微软雅黑", anchor=MSO_ANCHOR.TOP):
    """添加文本框"""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    tf.auto_size = None
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = font_name
    p.alignment = align
    tf.paragraphs[0].space_before = Pt(0)
    tf.paragraphs[0].space_after = Pt(0)
    return txBox


def add_multi_text(slide, left, top, width, height, lines, font_size=16,
                   color=PPTX_DARK, line_spacing=1.5, font_name="微软雅黑",
                   align=PP_ALIGN.LEFT):
    """添加多行文本"""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, line in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = line
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.name = font_name
        p.alignment = align
        p.space_after = Pt(font_size * (line_spacing - 1))
    return txBox


def add_page_number(slide, num, total=18):
    """添加页码"""
    add_text_box(slide, Inches(12.0), Inches(7.0), Inches(1.2), Inches(0.4),
                 f"{num:02d} / {total:02d}", font_size=10,
                 color=PPTX_WARM_GRAY, align=PP_ALIGN.RIGHT)


def add_section_title(slide, title, subtitle=None):
    """标准内页标题区"""
    add_gold_bar(slide, Inches(0.6), Inches(0.5), Inches(1.2), Pt(3))
    add_text_box(slide, Inches(0.6), Inches(0.65), Inches(10), Inches(0.7),
                 title, font_size=30, bold=True, color=PPTX_NAVY)
    if subtitle:
        add_text_box(slide, Inches(0.6), Inches(1.25), Inches(10), Inches(0.4),
                     subtitle, font_size=14, color=PPTX_WARM_GRAY)


def add_image_right(slide, img_path, width=Inches(5.5)):
    """右侧图片"""
    img_h = int(SLIDE_H) - Inches(1.0)
    left = int(SLIDE_W) - width - Inches(0.5)
    slide.shapes.add_picture(img_path, left, Inches(0.5), width=width)


def add_image_bottom(slide, img_path, height=Inches(4.2)):
    """底部图片"""
    img_w = int(SLIDE_W) - Inches(1.2)
    slide.shapes.add_picture(img_path, Inches(0.6), Inches(2.5), width=img_w)


# ── 构建每一页 ──

# ===================== 1. 封面 =====================
slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
add_shape_bg(slide, PPTX_NAVY)

# 封面中央金色装饰线
add_gold_bar(slide, Inches(4.5), Inches(1.8), Inches(4.3), Pt(2))
add_text_box(slide, Inches(1.5), Inches(2.0), Inches(10.3), Inches(1.2),
             "会议服务手册", font_size=52, bold=True,
             color=PPTX_GOLD, align=PP_ALIGN.CENTER)
add_text_box(slide, Inches(1.5), Inches(3.2), Inches(10.3), Inches(0.6),
             "CONFERENCE & MEETING SERVICE MANUAL", font_size=18,
             color=PPTX_LIGHT_GOLD, align=PP_ALIGN.CENTER)
add_gold_bar(slide, Inches(4.5), Inches(3.9), Inches(4.3), Pt(2))

add_text_box(slide, Inches(1.5), Inches(4.3), Inches(10.3), Inches(0.5),
             "尊享会议体验 · 专业服务保障", font_size=22,
             color=PPTX_WHITE, align=PP_ALIGN.CENTER)

add_text_box(slide, Inches(1.5), Inches(6.3), Inches(10.3), Inches(0.4),
             "[酒店名称]  |  高端会议服务品牌", font_size=14,
             color=PPTX_WARM_GRAY, align=PP_ALIGN.CENTER)

# ===================== 2. 目录 =====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_section_title(slide, "目 录", "CONTENTS")

toc_items = [
    ("01", "酒店简介", "酒店概况与会议中心"),
    ("02", "服务理念", "尊享服务承诺"),
    ("03", "会议场地", "各规格会议空间"),
    ("04", "设施设备", "专业会议配套"),
    ("05", "会议餐饮", "茶歇与宴会服务"),
    ("06", "商务服务", "商务中心支持"),
    ("07", "会议套餐", "灵活方案选择"),
    ("08", "服务流程", "一站式会议管理"),
    ("09", "专属管家", "一对一会议管家"),
    ("10", "住宿配套", "客房与VIP接待"),
    ("11", "交通便利", "出行无忧"),
    ("12", "成功案例", "知名企业信赖"),
    ("13", "联系我们", "即刻预订"),
]

x_start = Inches(0.8)
y_start = Inches(1.8)
col_w = Inches(5.8)
row_h = Inches(0.48)
for i, (num, title, desc) in enumerate(toc_items):
    col = i // 7
    row = i % 7
    x = x_start + col * col_w
    y = y_start + row * row_h
    # 序号
    add_text_box(slide, x, y, Inches(0.5), row_h,
                 num, font_size=14, bold=True, color=PPTX_GOLD, align=PP_ALIGN.LEFT)
    # 竖线分隔
    add_gold_bar(slide, x + Inches(0.55), y + Pt(4), Pt(1), Inches(0.28))
    # 标题
    add_text_box(slide, x + Inches(0.7), y, Inches(1.8), row_h,
                 title, font_size=15, bold=True, color=PPTX_DARK)
    # 描述
    add_text_box(slide, x + Inches(2.6), y + Pt(2), Inches(3.0), row_h,
                 desc, font_size=11, color=PPTX_WARM_GRAY)


# ===================== 3. 酒店简介 =====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_section_title(slide, "酒店简介", "ABOUT THE HOTEL")

intro_text = [
    "[酒店名称]坐落于城市核心地段，是一家集商务、会议、休闲于一体",
    "的高端奢华酒店。拥有超过3,000平方米的会议及宴会空间，配备",
    "世界一流的视听设备和专业的会议服务团队。",
    "",
    "我们秉承「奢华于形，尊贵于心」的服务理念，致力于为每一位宾客",
    "打造难忘的会议体验。累计服务超过500场国际会议、企业年会、",
    "商务洽谈及高端社交活动。",
]
add_multi_text(slide, Inches(0.8), Inches(2.0), Inches(7.0), Inches(4.0),
               intro_text, font_size=15, color=PPTX_DARK, line_spacing=1.6)

# 右侧数据卡片
data_items = [
    ("3,000+", "会议面积(m²)"),
    ("500+", "成功会议"),
    ("12", "多功能厅"),
    ("1,200", "最大容纳人数"),
]
for i, (num, label) in enumerate(data_items):
    x = Inches(8.5)
    y = Inches(2.0 + i * 1.3)
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, Inches(4.2), Inches(1.0))
    shape.fill.solid()
    shape.fill.fore_color.rgb = PPTX_NAVY
    shape.line.fill.background()
    add_text_box(slide, x + Inches(0.3), y + Pt(8), Inches(3.6), Inches(0.5),
                 num, font_size=28, bold=True, color=PPTX_GOLD, align=PP_ALIGN.CENTER)
    add_text_box(slide, x + Inches(0.3), y + Inches(0.5), Inches(3.6), Inches(0.35),
                 label, font_size=12, color=PPTX_WHITE, align=PP_ALIGN.CENTER)

# ===================== 4. 服务理念 =====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_section_title(slide, "服务理念", "SERVICE PHILOSOPHY")

pillars = [
    ("专业保障", "资深会议管家全程跟进，确保每个环节精准无误"),
    ("定制服务", "根据会议性质提供个性化布场与流程设计"),
    ("高效响应", "15分钟响应机制，会议期间问题即时解决"),
    ("奢华体验", "五星级酒店标准，从茶歇到宴请尽显品味"),
]
for i, (title, desc) in enumerate(pillars):
    x = Inches(0.8 + i * 3.1)
    y = Inches(2.2)
    # 卡片背景
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                    x, y, Inches(2.8), Inches(4.0))
    shape.fill.solid()
    shape.fill.fore_color.rgb = PPTX_NAVY
    shape.line.fill.background()
    # 金色顶部装饰
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                  x, y, Inches(2.8), Pt(4))
    bar.fill.solid()
    bar.fill.fore_color.rgb = PPTX_GOLD
    bar.line.fill.background()
    # 标题
    add_text_box(slide, x + Inches(0.2), y + Inches(0.4), Inches(2.4), Inches(0.5),
                 title, font_size=22, bold=True, color=PPTX_GOLD, align=PP_ALIGN.CENTER)
    # 描述
    add_text_box(slide, x + Inches(0.2), y + Inches(1.2), Inches(2.4), Inches(2.5),
                 desc, font_size=14, color=PPTX_WHITE, align=PP_ALIGN.CENTER)

# ===================== 5. 会议场地 =====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_section_title(slide, "会议场地", "MEETING VENUES")

venues = [
    ("大宴会厅", "1,200m²", "800人", "豪华宴会、大型峰会、年会盛典"),
    ("多功能厅 A", "500m²", "300人", "中型会议、产品发布会"),
    ("多功能厅 B", "300m²", "180人", "商务会议、培训讲座"),
    ("董事会议室", "120m²", "30人", "高层会议、私密洽谈"),
    ("VIP会客厅", "80m²", "20人", "贵宾接待、签约仪式"),
]
# 表格标题
headers = ["场地名称", "面积", "容纳人数", "适用场景"]
col_widths = [Inches(2.5), Inches(1.5), Inches(1.5), Inches(4.5)]
x_positions = [Inches(0.8)]
for w in col_widths[:-1]:
    x_positions.append(x_positions[-1] + w)

# 表头
for j, (hdr, xp, cw) in enumerate(zip(headers, x_positions, col_widths)):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                    xp, Inches(2.0), cw, Inches(0.5))
    shape.fill.solid()
    shape.fill.fore_color.rgb = PPTX_NAVY
    shape.line.color.rgb = PPTX_GOLD
    shape.line.width = Pt(0.5)
    add_text_box(slide, xp, Inches(2.05), cw, Inches(0.4),
                 hdr, font_size=13, bold=True, color=PPTX_GOLD, align=PP_ALIGN.CENTER)

# 数据行
for i, row_data in enumerate(venues):
    y = Inches(2.55 + i * 0.55)
    for j, (cell, xp, cw) in enumerate(zip(row_data, x_positions, col_widths)):
        shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                        xp, y, cw, Inches(0.5))
        shape.fill.solid()
        shape.fill.fore_color.rgb = PPTX_CREAM if i % 2 == 0 else PPTX_WHITE
        shape.line.color.rgb = PPTX_LIGHT_GOLD
        shape.line.width = Pt(0.5)
        fc = PPTX_GOLD if j == 0 else PPTX_DARK
        fb = True if j == 0 else False
        add_text_box(slide, xp, y + Pt(4), cw, Inches(0.4),
                     cell, font_size=12, bold=fb, color=fc, align=PP_ALIGN.CENTER)

# ===================== 6. 大宴会厅（带图）=====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_section_title(slide, "大宴会厅", "GRAND BALLROOM — 盛会之境")

banquet_img = gen_banquet_hall()
add_image_right(slide, banquet_img, Inches(6.0))

ballroom_info = [
    "• 面积：1,200m²（无柱设计）",
    "• 层高：8.5米",
    "• 容纳：800人剧院式 / 500人宴会式",
    "• 配备：LED大屏、顶尖灯光音响系统",
    "• 可分割为三个独立区域",
    "• 独立宴会厨房，确保菜品品质",
    "• 专属VIP休息室及化妆间",
]
add_multi_text(slide, Inches(0.8), Inches(2.0), Inches(5.5), Inches(4.5),
               ballroom_info, font_size=14, color=PPTX_DARK, line_spacing=1.7)

# ===================== 7. 中小会议室（带图）=====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_section_title(slide, "中小型会议室", "FLEXIBLE MEETING ROOMS")

room_img = gen_meeting_room()
add_image_right(slide, room_img, Inches(6.0))

room_info = [
    "• 多功能厅：500m² / 300m²，灵活分隔",
    "• 董事会议室：120m²，尊享私密空间",
    "• 全系智能照明与电动窗帘系统",
    "• 内置4K投影仪与电动幕布",
    "• 会议讨论系统（投票/同传）",
    "• 可移动隔断，多种布局自由切换",
    "• 会议礼宾全程值守",
]
add_multi_text(slide, Inches(0.8), Inches(2.0), Inches(5.5), Inches(4.5),
               room_info, font_size=14, color=PPTX_DARK, line_spacing=1.7)

# ===================== 8. 设施设备 =====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_section_title(slide, "会议设施与设备", "FACILITIES & EQUIPMENT")

equip_items = [
    ("🖥️ 视听系统", "LED大屏 / 投影仪 / 激光翻页笔"),
    ("🎤 音频系统", "专业音响 / 无线麦克 / 同声传译"),
    ("🌐 网络通讯", "千兆光纤WiFi / 视频会议系统"),
    ("📹 录制直播", "多机位录制 / 实时直播推流"),
    ("📋 会议配套", "电子签到 / 投票系统 / 电子席卡"),
    ("🪑 会场布置", "多场景布场 / 舞台搭建 / 灯光设计"),
]
for i, (title, desc) in enumerate(equip_items):
    col = i % 3
    row = i // 3
    x = Inches(0.8 + col * 4.1)
    y = Inches(2.2 + row * 2.2)
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                    x, y, Inches(3.8), Inches(1.8))
    shape.fill.solid()
    shape.fill.fore_color.rgb = PPTX_NAVY
    shape.line.fill.background()
    add_text_box(slide, x + Inches(0.3), y + Inches(0.3), Inches(3.2), Inches(0.5),
                 title, font_size=18, bold=True, color=PPTX_GOLD)
    add_text_box(slide, x + Inches(0.3), y + Inches(1.0), Inches(3.2), Inches(0.6),
                 desc, font_size=13, color=PPTX_WHITE)

# ===================== 9. 会议餐饮（带图）=====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_section_title(slide, "会议餐饮服务", "CATERING & REFRESHMENTS")

tea_img = gen_tea_break()
add_image_right(slide, tea_img, Inches(6.0))

catering_info = [
    "• 精选茶歇：高端咖啡茶饮、精致糕点水果",
    "• 商务套餐：中西式自助午餐、商务简餐",
    "• 主题晚宴：定制宴会菜单、酒会/冷餐会",
    "• 特殊需求：清真/素食/过敏源餐食说明",
    "• 服务标准：30分钟快速布场，实时补给",
    "• 酒水服务：精选红酒、香槟及特调饮品",
]
add_multi_text(slide, Inches(0.8), Inches(2.0), Inches(5.5), Inches(4.5),
               catering_info, font_size=14, color=PPTX_DARK, line_spacing=1.7)

# ===================== 10. 商务服务（带图）=====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_section_title(slide, "商务中心服务", "BUSINESS CENTER")

biz_img = gen_business_center()
add_image_right(slide, biz_img, Inches(6.0))

biz_info = [
    "• 文印服务：打印/复印/扫描/装订",
    "• 秘书服务：会议记录、文件翻译",
    "• 设备租用：笔记本电脑、平板设备",
    "• 快递服务：文件/物品寄送代收",
    "• 商务支持：名片制作、席卡打印",
    "• 24小时开放，随时满足商务需求",
]
add_multi_text(slide, Inches(0.8), Inches(2.0), Inches(5.5), Inches(4.5),
               biz_info, font_size=14, color=PPTX_DARK, line_spacing=1.7)

# ===================== 11. 会议套餐 =====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_section_title(slide, "会议套餐方案", "MEETING PACKAGES")

packages = [
    ("标准会议套餐", "¥288/位", [
        "会议室半天使用",
        "标准投影/音响设备",
        "上下午各一次茶歇",
        "商务午餐一份",
        "会议文具套装",
    ]),
    ("商务会议套餐", "¥488/位", [
        "会议室全天使用",
        "LED大屏/专业音响",
        "精致茶歇（无限续）",
        "自助午餐+下午茶",
        "全程会议管家服务",
    ]),
    ("尊享会议套餐", "¥888/位", [
        "全天使用+VIP会客厅",
        "全套视听+同声传译",
        "奢华茶歇+定制晚宴",
        "VIP客房一间一晚",
        "专车接送服务",
    ]),
]

for i, (name, price, items) in enumerate(packages):
    x = Inches(0.8 + i * 4.1)
    y = Inches(2.0)
    # 套餐卡片
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                    x, y, Inches(3.8), Inches(5.0))
    shape.fill.solid()
    shape.fill.fore_color.rgb = PPTX_NAVY
    shape.line.fill.background()
    # 顶部金色
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                  x, y, Inches(3.8), Pt(5))
    bar.fill.solid()
    bar.fill.fore_color.rgb = PPTX_GOLD
    bar.line.fill.background()
    # 套餐名
    add_text_box(slide, x + Inches(0.2), y + Inches(0.3), Inches(3.4), Inches(0.4),
                 name, font_size=18, bold=True, color=PPTX_GOLD, align=PP_ALIGN.CENTER)
    # 价格
    add_text_box(slide, x + Inches(0.2), y + Inches(0.8), Inches(3.4), Inches(0.5),
                 price, font_size=26, bold=True, color=PPTX_WHITE, align=PP_ALIGN.CENTER)
    # 分隔线
    add_gold_bar(slide, x + Inches(0.4), y + Inches(1.4), Inches(3.0), Pt(1))
    # 项目列表
    items_text = "\n".join([f"✓ {it}" for it in items])
    add_text_box(slide, x + Inches(0.3), y + Inches(1.6), Inches(3.2), Inches(3.0),
                 items_text, font_size=13, color=PPTX_LIGHT_GOLD, align=PP_ALIGN.LEFT)

# ===================== 12. 服务流程 =====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_section_title(slide, "会议服务流程", "SERVICE PROCESS")

steps = [
    ("01", "需求沟通", "了解会议性质、规模、预算"),
    ("02", "方案定制", "出具个性化会议方案与报价"),
    ("03", "签订合同", "确认服务细节，锁定档期"),
    ("04", "会前准备", "布场、设备调试、菜单确认"),
    ("05", "会中服务", "全程管家跟场，实时解决问题"),
    ("06", "会后跟进", "结算、反馈收集、客户档案归档"),
]

for i, (num, title, desc) in enumerate(steps):
    x = Inches(1.0 + i * 2.0)
    y = Inches(2.5)
    # 连线
    if i < len(steps) - 1:
        line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                       x + Inches(1.3), y + Inches(0.7),
                                       Inches(0.9), Pt(2))
        line.fill.solid()
        line.fill.fore_color.rgb = PPTX_GOLD
        line.line.fill.background()
    # 圆形序号
    shape = slide.shapes.add_shape(MSO_SHAPE.OVAL,
                                    x + Inches(0.3), y, Inches(0.8), Inches(0.8))
    shape.fill.solid()
    shape.fill.fore_color.rgb = PPTX_NAVY
    shape.line.color.rgb = PPTX_GOLD
    shape.line.width = Pt(2)
    add_text_box(slide, x + Inches(0.3), y + Pt(8), Inches(0.8), Inches(0.6),
                 num, font_size=20, bold=True, color=PPTX_GOLD, align=PP_ALIGN.CENTER)
    # 标题
    add_text_box(slide, x, y + Inches(1.0), Inches(1.8), Inches(0.4),
                 title, font_size=15, bold=True, color=PPTX_DARK, align=PP_ALIGN.CENTER)
    # 描述
    add_text_box(slide, x, y + Inches(1.4), Inches(1.8), Inches(0.6),
                 desc, font_size=11, color=PPTX_WARM_GRAY, align=PP_ALIGN.CENTER)
    add_page_number(slide, 12)

# ===================== 13. 专属管家（带图）=====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_section_title(slide, "专属会议管家", "PERSONAL EVENT MANAGER")

concierge_img = gen_concierge()
add_image_right(slide, concierge_img, Inches(6.0))

butler_info = [
    "每一位会议客户均配备专属会议管家，",
    "从会前筹备到现场执行提供一对一的全流程服务。",
    "",
    "✦ 会前：需求分析、方案策划、供应商协调",
    "✦ 会中：现场统筹、进度把控、应急处理",
    "✦ 会后：费用结算、效果评估、档案归档",
    "",
    "会议管家均经过国际金钥匙组织认证培训，",
    "精通双语服务，确保国际会议无缝对接。",
]
add_multi_text(slide, Inches(0.8), Inches(2.0), Inches(5.5), Inches(4.5),
               butler_info, font_size=14, color=PPTX_DARK, line_spacing=1.6)

# ===================== 14. 住宿配套（带图）=====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_section_title(slide, "住宿配套", "ACCOMMODATION")

vip_img = gen_vip_lounge()
add_image_right(slide, vip_img, Inches(6.0))

accomm_info = [
    "• 豪华客房与套房共计388间",
    "• 行政酒廊：专属楼层，尊享早餐与Happy Hour",
    "• VIP接待：总统套房专享会客区",
    "• 会议专属房价优惠，灵活入住/退房",
    "• 客房配备：高端床品、智能控制系统",
    "• 会议团队享免费健身房与室内泳池",
]
add_multi_text(slide, Inches(0.8), Inches(2.0), Inches(5.5), Inches(4.5),
               accomm_info, font_size=14, color=PPTX_DARK, line_spacing=1.7)

# ===================== 15. 交通便利 =====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_section_title(slide, "交通与出行", "TRANSPORTATION")

transport_info = [
    ("✈️ 机场接送", "提供机场VIP通道接机服务\n商务专车全程接送\n距国际机场仅20分钟车程"),
    ("🚇 地铁交通", "地铁1号线、3号线交汇\n步行至地铁站仅3分钟\n城市核心交通枢纽"),
    ("🚗 自驾停车", "300个地下停车位\n专设会议VIP停车区\n免费停车服务"),
    ("🚌 团队接送", "50座豪华大巴可直达\n提供团队专属接送方案\n多站点灵活停靠"),
]

for i, (title, desc) in enumerate(transport_info):
    col = i % 2
    row = i // 2
    x = Inches(0.8 + col * 6.2)
    y = Inches(2.0 + row * 2.5)
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                    x, y, Inches(5.8), Inches(2.2))
    shape.fill.solid()
    shape.fill.fore_color.rgb = PPTX_NAVY
    shape.line.fill.background()
    add_text_box(slide, x + Inches(0.3), y + Inches(0.2), Inches(5.2), Inches(0.4),
                 title, font_size=18, bold=True, color=PPTX_GOLD)
    add_text_box(slide, x + Inches(0.3), y + Inches(0.8), Inches(5.2), Inches(1.2),
                 desc, font_size=13, color=PPTX_WHITE)

# ===================== 16. 成功案例（带图）=====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_section_title(slide, "成功案例", "SUCCESS STORIES")

case_img = gen_success_case()
add_image_right(slide, case_img, Inches(6.0))

cases = [
    "✦ 国际科技峰会（800人）—— 连续三年指定合作酒店",
    "✦ 全球金融论坛（500人）—— 同声传译6语种全程保障",
    "✦ 年度经销商大会（400人）—— 一站式会务托管服务",
    "✦ 高端品牌发布会（200人）—— 灯光舞美定制方案",
    "✦ 跨国企业董事会（30人）—— 总统套房私密商务洽谈",
]
add_multi_text(slide, Inches(0.8), Inches(2.0), Inches(5.5), Inches(4.5),
               cases, font_size=14, color=PPTX_DARK, line_spacing=1.8)

# ===================== 17. 联系我们 =====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, PPTX_NAVY)

add_gold_bar(slide, Inches(5.0), Inches(1.0), Inches(3.3), Pt(2))
add_text_box(slide, Inches(1.0), Inches(1.3), Inches(11.3), Inches(0.8),
             "联系我们", font_size=40, bold=True,
             color=PPTX_GOLD, align=PP_ALIGN.CENTER)
add_text_box(slide, Inches(1.0), Inches(2.1), Inches(11.3), Inches(0.5),
             "即刻预订，尊享非凡会议体验", font_size=20,
             color=PPTX_WHITE, align=PP_ALIGN.CENTER)
add_gold_bar(slide, Inches(5.0), Inches(2.7), Inches(3.3), Pt(2))

contact_items = [
    ("📍 地址", "城市中心CBD · [详细地址]"),
    ("📞 预订热线", "400-XXX-XXXX"),
    ("📧 电子邮箱", "conference@[hotel].com"),
    ("🌐 官方网站", "www.[hotel].com/meeting"),
    ("⏰ 服务时间", "24小时会议服务专线"),
]
for i, (icon, info) in enumerate(contact_items):
    y = Inches(3.2 + i * 0.65)
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                    Inches(3.5), y, Inches(6.3), Inches(0.5))
    shape.fill.solid()
    shape.fill.fore_color.rgb = PPTX_DARK
    shape.line.color.rgb = PPTX_GOLD
    shape.line.width = Pt(1)
    add_text_box(slide, Inches(3.7), y + Pt(4), Inches(5.9), Inches(0.4),
                 f"{icon}  {info}", font_size=15, color=PPTX_LIGHT_GOLD, align=PP_ALIGN.CENTER)

add_text_box(slide, Inches(1.0), Inches(6.8), Inches(11.3), Inches(0.4),
             "扫描二维码 · 在线预订 · 即时报价", font_size=12,
             color=PPTX_WARM_GRAY, align=PP_ALIGN.CENTER)

# ===================== 18. 封底 =====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_shape_bg(slide, PPTX_NAVY)

add_gold_bar(slide, Inches(5.0), Inches(2.5), Inches(3.3), Pt(2))
add_text_box(slide, Inches(1.0), Inches(2.8), Inches(11.3), Inches(0.8),
             "感谢您的信赖", font_size=40, bold=True,
             color=PPTX_GOLD, align=PP_ALIGN.CENTER)
add_text_box(slide, Inches(1.0), Inches(3.7), Inches(11.3), Inches(0.5),
             "以专业致敬每一场重要的会议", font_size=20,
             color=PPTX_WHITE, align=PP_ALIGN.CENTER)
add_gold_bar(slide, Inches(5.0), Inches(4.3), Inches(3.3), Pt(2))

add_text_box(slide, Inches(1.0), Inches(5.5), Inches(11.3), Inches(0.4),
             "[酒店名称]  |  高端会议服务品牌", font_size=14,
             color=PPTX_WARM_GRAY, align=PP_ALIGN.CENTER)

# ── 保存 ──
output_path = os.path.join(PPTX_DIR, "酒店会议服务手册.pptx")
prs.save(output_path)
print(f"[OK] PPT已生成: {output_path}")
print(f"[OK] 配图已生成到: {IMG_DIR}")
