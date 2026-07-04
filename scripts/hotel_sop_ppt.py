#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""星级酒店会议服务 SOP 手册 PPT — 实操版"""

import os, math
from PIL import Image, ImageDraw, ImageFont
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

OUTPUT_DIR = os.path.join(os.path.dirname(os.path.abspath(__file__)), "..", "outputs")
PPTX_DIR = os.path.join(OUTPUT_DIR, "pptx")
IMG_DIR = os.path.join(OUTPUT_DIR, "images")
os.makedirs(PPTX_DIR, exist_ok=True)
os.makedirs(IMG_DIR, exist_ok=True)

# ── Colors ──
GOLD = (180, 150, 85)
NAVY = (25, 30, 50)
DARK = (40, 40, 50)
WHITE = (255, 255, 255)
CREAM = (248, 245, 240)
LIGHT_GRAY = (230, 228, 222)
WARM_GRAY = (160, 155, 145)
TABLE_COLOR = (60, 65, 85)
CHAIR_COLOR = (140, 130, 110)
CHAIR_BRIGHT = (180, 170, 150)
STAGE_COLOR = (100, 95, 85)

PPTX_GOLD = RGBColor(0xB4, 0x96, 0x55)
PPTX_NAVY = RGBColor(0x19, 0x1E, 0x32)
PPTX_DARK = RGBColor(0x28, 0x28, 0x32)
PPTX_WHITE = RGBColor(0xFF, 0xFF, 0xFF)
PPTX_CREAM = RGBColor(0xF8, 0xF5, 0xF0)
PPTX_LIGHT_GRAY = RGBColor(0xE6, 0xE4, 0xDE)
PPTX_WARM_GRAY = RGBColor(0xA0, 0x9B, 0x91)
PPTX_ACCENT = RGBColor(0xD4, 0xA0, 0x40)

def find_font(size=20, bold=False):
    candidates = [
        "C:/Windows/Fonts/msyhbd.ttc",
        "C:/Windows/Fonts/simhei.ttf",
        "C:/Windows/Fonts/msyh.ttc",
        "C:/Windows/Fonts/simsun.ttc",
    ]
    if not bold:
        candidates = [c for c in candidates if 'bd' not in c.lower()] + candidates
    for p in candidates:
        if os.path.exists(p):
            try:
                return ImageFont.truetype(p, size)
            except:
                continue
    return ImageFont.load_default()

# ═══════════════════════════════════════════
#  TOP-DOWN VENUE LAYOUT DRAWING FUNCTIONS
# ═══════════════════════════════════════════

def draw_stage(draw, cx, y, w, label="讲台"):
    """讲台"""
    draw.rectangle([cx - w//2, y - 15, cx + w//2, y + 15], fill=STAGE_COLOR)
    draw.rectangle([cx - w//2, y - 15, cx + w//2, y - 12], fill=GOLD)
    f = find_font(14)
    draw.text((cx - 15, y - 10), label, fill=WHITE, font=f)

def draw_rect_table(draw, x, y, w, h, fill=None, outline=None):
    """矩形桌"""
    fill = fill or TABLE_COLOR
    outline = outline or GOLD
    draw.rectangle([x, y, x + w, y + h], fill=fill, outline=outline, width=2)

def draw_round_table(draw, cx, cy, r, fill=None, outline=None):
    """圆桌"""
    fill = fill or TABLE_COLOR
    outline = outline or GOLD
    draw.ellipse([cx - r, cy - r, cx + r, cy + r], fill=fill, outline=outline, width=2)

def draw_chair(draw, cx, cy, angle=0, color=None):
    """椅子（小矩形）"""
    color = color or CHAIR_COLOR
    w, h = 10, 14
    rad = math.radians(angle or 0)
    # 围绕桌子摆放
    draw.ellipse([cx - w//2, cy - h//2, cx + w//2, cy + h//2], fill=color, outline=GOLD, width=1)

def draw_chairs_around_rect(draw, tx, ty, tw, th, spacing=22, side_flags="tbrl"):
    """矩形桌周围摆椅子"""
    # top
    if 't' in side_flags:
        for x in range(tx + 12, tx + tw - 8, spacing):
            draw_chair(draw, x, ty - 8, 0)
    # bottom
    if 'b' in side_flags:
        for x in range(tx + 12, tx + tw - 8, spacing):
            draw_chair(draw, x, ty + th + 8, 180)
    # left
    if 'l' in side_flags:
        for y in range(ty + 12, ty + th - 8, spacing):
            draw_chair(draw, tx - 8, y, 270)
    # right
    if 'r' in side_flags:
        for y in range(ty + 12, ty + th - 8, spacing):
            draw_chair(draw, tx + tw + 8, y, 90)

def draw_grid_lines(draw, w, h):
    """辅助网格（极淡）"""
    for x in range(0, w, 40):
        draw.line([(x, 0), (x, h)], fill=(220, 218, 212), width=1)
    for y in range(0, h, 40):
        draw.line([(0, y), (w, y)], fill=(220, 218, 212), width=1)

# ── Layout Generators ──

W, H = 1200, 780

def gen_theater():
    """剧院式"""
    img = Image.new("RGB", (W, H), CREAM)
    draw = ImageDraw.Draw(img)
    draw_grid_lines(draw, W, H)
    draw_stage(draw, W//2, 60, 200)
    rows, cols = 10, 14
    start_x, start_y = 60, 120
    gap_x = (W - 120) // cols
    gap_y = 55
    chair_w, chair_h = 12, 16
    for r in range(rows):
        for c in range(cols):
            x = start_x + c * gap_x
            y = start_y + r * gap_y
            # 走道
            if c == cols // 2:
                continue
            draw.ellipse([x - chair_w//2, y - chair_h//2, x + chair_w//2, y + chair_h//2],
                         fill=CHAIR_COLOR, outline=GOLD, width=1)
    # 标注
    f = find_font(22, bold=True)
    draw.text((W//2 - 100, H - 50), "剧院式 THEATER", fill=DARK, font=f)
    img.save(os.path.join(IMG_DIR, "sop_theater.png"))
    return os.path.join(IMG_DIR, "sop_theater.png")

def gen_classroom():
    """课桌式"""
    img = Image.new("RGB", (W, H), CREAM)
    draw = ImageDraw.Draw(img)
    draw_grid_lines(draw, W, H)
    draw_stage(draw, W//2, 60, 200)
    table_w, table_h = 70, 45
    gap_x = 90
    gap_y = 80
    cols = 10
    rows = 7
    offset_x = (W - cols * gap_x + gap_x) // 2
    start_y = 130
    for r in range(rows):
        for c in range(cols):
            x = offset_x + c * gap_x
            y = start_y + r * gap_y
            draw_rect_table(draw, x, y, table_w, table_h)
            draw_chairs_around_rect(draw, x, y, table_w, table_h, 18, "t")
    f = find_font(22, bold=True)
    draw.text((W//2 - 100, H - 50), "课桌式 CLASSROOM", fill=DARK, font=f)
    img.save(os.path.join(IMG_DIR, "sop_classroom.png"))
    return os.path.join(IMG_DIR, "sop_classroom.png")

def gen_ushape():
    """U型"""
    img = Image.new("RGB", (W, H), CREAM)
    draw = ImageDraw.Draw(img)
    draw_grid_lines(draw, W, H)
    draw_stage(draw, W//2, 50, 180)
    tw, th = 60, 35
    gap = 65
    # U shape coordinates
    cx = W // 2
    top_y = 100
    bottom_y = H - 100
    left_x = 100
    right_x = W - 100 - tw
    # Top row
    for i in range(-5, 6):
        x = cx + i * gap - tw//2
        if abs(x - cx) < 20 and i != 0:
            continue
        draw_rect_table(draw, x, top_y, tw, th)
        draw_chairs_around_rect(draw, x, top_y, tw, th, 18, "b")
    # Left column
    for i in range(1, 7):
        y = top_y + i * gap
        if y > bottom_y:
            break
        draw_rect_table(draw, left_x, y, tw, th)
        draw_chairs_around_rect(draw, left_x, y, tw, th, 18, "r")
    # Right column
    for i in range(1, 7):
        y = top_y + i * gap
        if y > bottom_y:
            break
        draw_rect_table(draw, right_x, y, tw, th)
        draw_chairs_around_rect(draw, right_x, y, tw, th, 18, "l")
    f = find_font(22, bold=True)
    draw.text((W//2 - 80, H - 50), "U型 U-SHAPE", fill=DARK, font=f)
    img.save(os.path.join(IMG_DIR, "sop_ushape.png"))
    return os.path.join(IMG_DIR, "sop_ushape.png")

def gen_hollow_square():
    """回字型"""
    img = Image.new("RGB", (W, H), CREAM)
    draw = ImageDraw.Draw(img)
    draw_grid_lines(draw, W, H)
    draw_stage(draw, W//2, 45, 160)
    tw, th = 60, 35
    gap = 65
    cx = W // 2
    cy = H // 2 + 20
    # 4 sides
    # Top
    for i in range(-5, 6):
        if abs(i * gap + cx) - cx < gap * 0.5:
            continue
        x = cx + i * gap - tw//2
        draw_rect_table(draw, x, cy - 130, tw, th)
        draw_chairs_around_rect(draw, x, cy - 130, tw, th, 18, "b")
    # Bottom
    for i in range(-5, 6):
        if abs(i * gap + cx) - cx < gap * 0.5:
            continue
        x = cx + i * gap - tw//2
        draw_rect_table(draw, x, cy + 100, tw, th)
        draw_chairs_around_rect(draw, x, cy + 100, tw, th, 18, "t")
    # Left
    for i in range(-2, 3):
        if i == 0:
            continue
        y = cy + i * gap - th//2
        draw_rect_table(draw, cx - 200, y, tw, th)
        draw_chairs_around_rect(draw, cx - 200, y, tw, th, 18, "r")
    # Right
    for i in range(-2, 3):
        if i == 0:
            continue
        y = cy + i * gap - th//2
        draw_rect_table(draw, cx + 140, y, tw, th)
        draw_chairs_around_rect(draw, cx + 140, y, tw, th, 18, "l")
    f = find_font(22, bold=True)
    draw.text((W//2 - 100, H - 50), "回字型 HOLLOW SQUARE", fill=DARK, font=f)
    img.save(os.path.join(IMG_DIR, "sop_hollow_square.png"))
    return os.path.join(IMG_DIR, "sop_hollow_square.png")

def gen_fishbone():
    """鱼骨式"""
    img = Image.new("RGB", (W, H), CREAM)
    draw = ImageDraw.Draw(img)
    draw_grid_lines(draw, W, H)
    draw_stage(draw, W//2, 40, 180)
    tw, th = 70, 40
    cx = W // 2
    cy = H // 2 + 30
    gap_x = 75
    gap_y = 75
    # Left side angled
    for i in range(1, 6):
        x = cx - 150 - i * gap_x
        y = cy - i * gap_y // 2
        # angled (diamond) by using offset
        draw_rect_table(draw, x, y, tw, th, fill=(55, 60, 75))
        draw.rectangle([x + 5, y + 5, x + tw - 5, y + th - 5], outline=GOLD + (50,), width=1)
        draw_chairs_around_rect(draw, x, y, tw, th, 18, "t")
        # second in row
        x2 = cx - 150 - i * gap_x
        y2 = cy + i * gap_y // 2
        draw_rect_table(draw, x2, y2, tw, th, fill=(55, 60, 75))
        draw.rectangle([x2 + 5, y2 + 5, x2 + tw - 5, y2 + th - 5], outline=GOLD + (50,), width=1)
        draw_chairs_around_rect(draw, x2, y2, tw, th, 18, "b")
    # Right side
    for i in range(1, 6):
        x = cx + 80 + i * gap_x
        y = cy - i * gap_y // 2
        draw_rect_table(draw, x, y, tw, th, fill=(55, 60, 75))
        draw.rectangle([x + 5, y + 5, x + tw - 5, y + th - 5], outline=GOLD + (50,), width=1)
        draw_chairs_around_rect(draw, x, y, tw, th, 18, "t")
        x2 = cx + 80 + i * gap_x
        y2 = cy + i * gap_y // 2
        draw_rect_table(draw, x2, y2, tw, th, fill=(55, 60, 75))
        draw.rectangle([x2 + 5, y2 + 5, x2 + tw - 5, y2 + th - 5], outline=GOLD + (50,), width=1)
        draw_chairs_around_rect(draw, x2, y2, tw, th, 18, "b")
    f = find_font(22, bold=True)
    draw.text((W//2 - 100, H - 50), "鱼骨式 FISHBONE", fill=DARK, font=f)
    img.save(os.path.join(IMG_DIR, "sop_fishbone.png"))
    return os.path.join(IMG_DIR, "sop_fishbone.png")

def gen_banquet():
    """宴会式（圆桌）"""
    img = Image.new("RGB", (W, H), CREAM)
    draw = ImageDraw.Draw(img)
    draw_grid_lines(draw, W, H)
    draw_stage(draw, W//2, 30, 140)
    cx, cy = W // 2, H // 2 + 30
    # 圆桌矩阵
    positions = [
        (cx - 180, cy - 130), (cx, cy - 130), (cx + 180, cy - 130),
        (cx - 180, cy), (cx, cy), (cx + 180, cy),
        (cx - 180, cy + 130), (cx, cy + 130), (cx + 180, cy + 130),
    ]
    r = 45
    for px, py in positions:
        draw_round_table(draw, px, py, r)
        # 周围椅子
        for angle in range(0, 360, 45):
            rad = math.radians(angle)
            cx2 = px + (r + 14) * math.cos(rad)
            cy2 = py + (r + 14) * math.sin(rad)
            draw_chair(draw, cx2, cy2, angle)
    f = find_font(22, bold=True)
    draw.text((W//2 - 100, H - 50), "宴会式 BANQUET", fill=DARK, font=f)
    img.save(os.path.join(IMG_DIR, "sop_banquet.png"))
    return os.path.join(IMG_DIR, "sop_banquet.png")

def gen_cocktail():
    """鸡尾酒会式"""
    img = Image.new("RGB", (W, H), CREAM)
    draw = ImageDraw.Draw(img)
    draw_grid_lines(draw, W, H)
    # 高脚圆桌分散
    positions = [
        (100, 100), (300, 80), (500, 120), (700, 90), (900, 110), (1100, 100),
        (150, 250), (350, 230), (550, 260), (750, 240), (950, 250), (1050, 270),
        (200, 400), (400, 380), (600, 410), (800, 390), (1000, 400),
        (120, 540), (320, 520), (520, 550), (720, 530), (920, 540), (1080, 560),
        (250, 670), (450, 650), (650, 670), (850, 660), (1050, 680),
    ]
    for px, py in positions:
        r = 14
        draw_round_table(draw, px, py, r, fill=(80, 75, 65))
    # 吧台
    draw.rectangle([W - 200, 50, W - 30, 130], fill=(55, 50, 42))
    draw.rectangle([W - 200, 50, W - 30, 53], fill=GOLD)
    f_lbl = find_font(16)
    draw.text((W - 170, 85), "BAR", fill=GOLD, font=f_lbl)
    f = find_font(22, bold=True)
    draw.text((W//2 - 120, H - 50), "鸡尾酒会式 COCKTAIL", fill=DARK, font=f)
    img.save(os.path.join(IMG_DIR, "sop_cocktail.png"))
    return os.path.join(IMG_DIR, "sop_cocktail.png")

def gen_boardroom():
    """董事会式"""
    img = Image.new("RGB", (W, H), CREAM)
    draw = ImageDraw.Draw(img)
    draw_grid_lines(draw, W, H)
    cx, cy = W // 2, H // 2 + 20
    # 大椭圆会议桌
    draw.ellipse([cx - 300, cy - 110, cx + 300, cy + 110], fill=TABLE_COLOR, outline=GOLD, width=3)
    draw.ellipse([cx - 295, cy - 105, cx + 295, cy + 105], fill=(50, 55, 70), outline=None)
    # 周围椅子
    for angle in range(0, 360, 20):
        rad = math.radians(angle)
        px = cx + 170 * math.cos(rad)
        py = cy + 65 * math.sin(rad)
        draw_chair(draw, px, py, angle, color=CHAIR_BRIGHT)
    f = find_font(22, bold=True)
    draw.text((W//2 - 120, H - 50), "董事会式 BOARDROOM", fill=DARK, font=f)
    img.save(os.path.join(IMG_DIR, "sop_boardroom.png"))
    return os.path.join(IMG_DIR, "sop_boardroom.png")

def gen_t_shape():
    """T型台"""
    img = Image.new("RGB", (W, H), CREAM)
    draw = ImageDraw.Draw(img)
    draw_grid_lines(draw, W, H)
    cx = W // 2
    # T台
    draw.rectangle([cx - 30, 40, cx + 30, 140], fill=STAGE_COLOR)
    draw.rectangle([cx - 30, 40, cx + 30, 43], fill=GOLD)
    # T台横向
    draw.rectangle([cx - 120, 80, cx + 120, 120], fill=STAGE_COLOR)
    draw.rectangle([cx - 120, 80, cx + 120, 83], fill=GOLD)
    # 竖排椅子
    for row in range(8):
        for col in range(12):
            x = 60 + col * 95
            y = 170 + row * 65
            draw.ellipse([x - 10, y - 10, x + 10, y + 10], fill=CHAIR_COLOR, outline=GOLD, width=1)
    f = find_font(22, bold=True)
    draw.text((W//2 - 80, H - 50), "T型台 T-SHAPE", fill=DARK, font=f)
    img.save(os.path.join(IMG_DIR, "sop_tshape.png"))
    return os.path.join(IMG_DIR, "sop_tshape.png")

def gen_equipment():
    """会议设备示意图"""
    img = Image.new("RGB", (W, H), CREAM)
    draw = ImageDraw.Draw(img)
    # 网格
    for x in range(0, W, 40):
        draw.line([(x, 0), (x, H)], fill=(225, 223, 218), width=1)
    for y in range(0, H, 40):
        draw.line([(0, y), (W, y)], fill=(225, 223, 218), width=1)

    # 投影仪
    cx, cy = 200, 200
    draw.rectangle([cx - 50, cy - 20, cx + 50, cy + 20], fill=(40, 40, 50))
    draw.ellipse([cx - 15, cy - 30, cx + 15, cy - 10], fill=(80, 80, 100))
    draw.text((cx - 30, cy - 60), "投影仪", fill=DARK, font=find_font(16))

    # 幕布
    cx2, cy2 = 600, 180
    draw.rectangle([cx2 - 120, cy2 - 60, cx2 + 120, cy2 + 60], fill=(200, 200, 210))
    draw.rectangle([cx2 - 120, cy2 - 60, cx2 + 120, cy2 - 57], fill=GOLD)
    draw.text((cx2 - 30, cy2 - 90), "投影幕布", fill=DARK, font=find_font(16))

    # 音响
    for i, (sx, sy) in enumerate([(100, 500), (1050, 500)]):
        draw.rectangle([sx - 30, sy - 60, sx + 30, sy + 60], fill=(30, 30, 40))
        draw.ellipse([sx - 20, sy - 20, sx + 20, sy + 20], fill=(50, 50, 60))
        draw.text((sx - 30, sy - 90), "音响", fill=DARK, font=find_font(16))

    # 讲台
    cx3 = 400
    draw.rectangle([cx3 - 60, 450, cx3 + 60, 520], fill=STAGE_COLOR)
    draw.rectangle([cx3 - 60, 450, cx3 + 60, 453], fill=GOLD)
    draw.text((cx3 - 30, 400), "讲台", fill=DARK, font=find_font(16))

    # 同声传译 booth
    cx4 = 900
    draw.rectangle([cx4 - 40, 450, cx4 + 40, 520], fill=(55, 60, 70))
    draw.rectangle([cx4 - 40, 450, cx4 + 40, 453], fill=GOLD)
    draw.text((cx4 - 60, 400), "同传间", fill=DARK, font=find_font(16))

    # 灯光
    for lx, ly in [(300, 80), (500, 60), (700, 80), (900, 60)]:
        draw.polygon([(lx, ly), (lx - 15, ly + 30), (lx + 15, ly + 30)], fill=GOLD + (150,))
    f = find_font(22, bold=True)
    draw.text((W//2 - 100, H - 50), "会议设备布局示意", fill=DARK, font=f)
    img.save(os.path.join(IMG_DIR, "sop_equipment.png"))
    return os.path.join(IMG_DIR, "sop_equipment.png")

def gen_tea_setup():
    """茶歇摆台"""
    img = Image.new("RGB", (W, H), CREAM)
    draw = ImageDraw.Draw(img)
    for x in range(0, W, 40):
        draw.line([(x, 0), (x, H)], fill=(225, 223, 218), width=1)
    for y in range(0, H, 40):
        draw.line([(0, y), (W, y)], fill=(225, 223, 218), width=1)
    cx = W // 2
    # 长台
    draw.rectangle([100, 250, W - 100, 370], fill=(65, 58, 48))
    draw.rectangle([100, 250, W - 100, 253], fill=GOLD)
    draw.rectangle([100, 367, W - 100, 370], fill=GOLD)
    # 饮品区
    draw.rectangle([120, 270, 350, 350], fill=(50, 45, 38))
    draw.text((180, 300), "☕ 咖啡/茶", fill=GOLD, font=find_font(20, bold=True))
    # 点心区
    draw.rectangle([370, 270, 650, 350], fill=(50, 45, 38))
    draw.text((420, 300), "🍰 甜品/水果", fill=GOLD, font=find_font(20, bold=True))
    # 冷餐区
    draw.rectangle([670, 270, 930, 350], fill=(50, 45, 38))
    draw.text((720, 300), "🥪 冷餐/沙拉", fill=GOLD, font=find_font(20, bold=True))
    # 餐具区
    draw.rectangle([950, 270, W - 120, 350], fill=(50, 45, 38))
    draw.text((1000, 300), "🍽️ 餐具纸巾", fill=GOLD, font=find_font(20, bold=True))
    f = find_font(22, bold=True)
    draw.text((W//2 - 100, H - 50), "茶歇摆台标准布局", fill=DARK, font=f)
    img.save(os.path.join(IMG_DIR, "sop_tea_setup.png"))
    return os.path.join(IMG_DIR, "sop_tea_setup.png")

# Generate all images
print("[生成配图中...]")
layout_images = {
    "theater": gen_theater(),
    "classroom": gen_classroom(),
    "ushape": gen_ushape(),
    "hollow_square": gen_hollow_square(),
    "fishbone": gen_fishbone(),
    "banquet": gen_banquet(),
    "cocktail": gen_cocktail(),
    "boardroom": gen_boardroom(),
    "tshape": gen_t_shape(),
    "equipment": gen_equipment(),
    "tea_setup": gen_tea_setup(),
}
print(f"[OK] 共生成 {len(layout_images)} 张配图")

# ═══════════════════════════════════════════
#  PPT GENERATION
# ═══════════════════════════════════════════

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
SW = prs.slide_width
SH = prs.slide_height

def add_bg(slide, color=PPTX_CREAM):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_shape_bg(slide, color=PPTX_NAVY):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SW, SH)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()

def add_gold_bar(slide, left, top, width, height=Pt(3)):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = PPTX_GOLD
    shape.line.fill.background()

def add_textbox(slide, left, top, width, height, text, size=18,
                bold=False, color=PPTX_DARK, align=PP_ALIGN.LEFT,
                font="微软雅黑"):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = font
    p.alignment = align
    return tb

def add_multiline(slide, left, top, width, height, lines, size=16,
                  color=PPTX_DARK, spacing=1.5, font="微软雅黑", align=PP_ALIGN.LEFT):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line
        p.font.size = Pt(size)
        p.font.color.rgb = color
        p.font.name = font
        p.alignment = align
        p.space_after = Pt(size * (spacing - 1) * 0.5)
    return tb

def add_section_title(slide, title, subtitle=None):
    bar = add_gold_bar(slide, Inches(0.5), Inches(0.45), Inches(1.0), Pt(3))
    add_textbox(slide, Inches(0.5), Inches(0.6), Inches(12), Inches(0.6),
                title, size=28, bold=True, color=PPTX_NAVY)
    if subtitle:
        add_textbox(slide, Inches(0.5), Inches(1.15), Inches(12), Inches(0.35),
                    subtitle, size=13, color=PPTX_WARM_GRAY)

def add_card(slide, left, top, width, height, title, body_lines,
             title_color=PPTX_GOLD, bg_color=PPTX_NAVY, title_size=18, body_size=13):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = bg_color
    shape.line.fill.background()
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, Pt(4))
    bar.fill.solid()
    bar.fill.fore_color.rgb = PPTX_GOLD
    bar.line.fill.background()
    add_textbox(slide, left + Inches(0.2), top + Inches(0.25), width - Inches(0.4), Inches(0.45),
                title, size=title_size, bold=True, color=title_color)
    add_multiline(slide, left + Inches(0.2), top + Inches(0.8), width - Inches(0.4), height - Inches(1.0),
                  body_lines, size=body_size, color=PPTX_WHITE, spacing=1.5)

def add_sop_table(slide, left, top, col_widths, headers, rows, header_bg=PPTX_NAVY, header_fg=PPTX_GOLD):
    x_positions = [left]
    for w in col_widths[:-1]:
        x_positions.append(x_positions[-1] + w)
    # Header
    for hdr, xp, cw in zip(headers, x_positions, col_widths):
        shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, xp, top, cw, Inches(0.4))
        shape.fill.solid()
        shape.fill.fore_color.rgb = header_bg
        shape.line.color.rgb = PPTX_GOLD
        shape.line.width = Pt(0.5)
        add_textbox(slide, xp, top + Pt(3), cw, Inches(0.35), hdr,
                    size=11, bold=True, color=header_fg, align=PP_ALIGN.CENTER)
    # Rows
    for i, row in enumerate(rows):
        y = top + Inches(0.4) + i * Inches(0.38)
        for j, (cell, xp, cw) in enumerate(zip(row, x_positions, col_widths)):
            shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, xp, y, cw, Inches(0.38))
            shape.fill.solid()
            shape.fill.fore_color.rgb = PPTX_CREAM if i % 2 == 0 else PPTX_WHITE
            shape.line.color.rgb = PPTX_LIGHT_GRAY
            shape.line.width = Pt(0.5)
            fc = PPTX_GOLD if j == 0 else PPTX_DARK
            fb = j == 0
            add_textbox(slide, xp, y + Pt(2), cw, Inches(0.34), cell,
                        size=10, bold=fb, color=fc, align=PP_ALIGN.CENTER)

def add_page_num(slide, num, total=22):
    add_textbox(slide, Inches(12.0), Inches(7.0), Inches(1.2), Inches(0.35),
                f"{num:02d}/{total:02d}", size=9, color=PPTX_WARM_GRAY, align=PP_ALIGN.RIGHT)

TOTAL_SLIDES = 22

# ══════════════════════════════════════
#  1. 封面
# ══════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_shape_bg(slide, PPTX_NAVY)
add_gold_bar(slide, Inches(4.8), Inches(1.5), Inches(3.7), Pt(2))
add_textbox(slide, Inches(1), Inches(1.7), Inches(11.3), Inches(1.0),
            "会议服务标准操作手册", size=46, bold=True, color=PPTX_GOLD, align=PP_ALIGN.CENTER)
add_textbox(slide, Inches(1), Inches(2.8), Inches(11.3), Inches(0.5),
            "CONFERENCE SERVICE SOP MANUAL", size=18, color=PPTX_LIGHT_GRAY, align=PP_ALIGN.CENTER)
add_gold_bar(slide, Inches(4.8), Inches(3.4), Inches(3.7), Pt(2))
add_multiline(slide, Inches(2), Inches(3.8), Inches(9.3), Inches(1.5),
              ["适用部门：宴会部 / 会议服务部 / 前厅部",
               "版本号：V1.0    实施日期：2026年1月",
               "密级：内部培训资料"],
              size=16, color=PPTX_WHITE, align=PP_ALIGN.CENTER, spacing=1.6)
add_textbox(slide, Inches(1), Inches(6.5), Inches(11.3), Inches(0.4),
            "星级酒店会议服务标准化流程", size=14, color=PPTX_WARM_GRAY, align=PP_ALIGN.CENTER)

# ══════════════════════════════════════
#  2. 目录
# ══════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_section_title(slide, "目  录", "CONTENTS")
toc = [
    ("01", "会议服务总则", "适用范围与职责划分"),
    ("02", "会议预订流程", "从接洽到确认的全流程"),
    ("03", "会前准备SOP", "场地布置与设备检查"),
    ("04", "剧院式布局", " Theater Style 布局标准+适用场景"),
    ("05", "课桌式布局", " Classroom Style 布局标准+适用场景"),
    ("06", "U型布局", " U-Shape Style 布局标准+适用场景"),
    ("07", "回字型布局", " Hollow Square 布局标准+适用场景"),
    ("08", "鱼骨式布局", " Fishbone Style 布局标准+适用场景"),
    ("09", "宴会式布局", " Banquet Style 布局标准+适用场景"),
    ("10", "鸡尾酒会式", " Cocktail Style 布局标准+适用场景"),
    ("11", "董事会式布局", " Boardroom Style 布局标准+适用场景"),
    ("12", "T型台布局", " T-Shape Style 布局标准+适用场景"),
    ("13", "台型选择指南", "各台型对比与推荐场景"),
    ("14", "会议设备标准", " AV设备清单与操作规范"),
    ("15", "会议摆台标准", "会议文具/水牌/席卡标准"),
    ("16", "茶歇服务SOP", "茶歇摆台与服务流程"),
    ("17", "会议餐饮服务", "宴会/自助餐/酒会服务标准"),
    ("18", "会议管家职责", "管家岗位职责与服务动线"),
    ("19", "应急预案", "常见突发情况处理方案"),
    ("20", "会议物品清单", "布场工具/文具/设备清单"),
    ("21", "附录", "标准表单/检查表/交接表"),
    ("22", "封底", ""),
]
x_start = Inches(0.6)
y_start = Inches(1.7)
col_w = Inches(6.0)
row_h = Inches(0.36)
for i, (num, title, desc) in enumerate(toc):
    col = i // 11
    row = i % 11
    x = x_start + col * col_w
    y = y_start + row * row_h
    add_textbox(slide, x, y, Inches(0.4), row_h, num, size=12, bold=True, color=PPTX_GOLD)
    add_gold_bar(slide, x + Inches(0.45), y + Pt(5), Pt(1), Inches(0.2))
    add_textbox(slide, x + Inches(0.6), y, Inches(2.2), row_h, title, size=13, bold=True, color=PPTX_DARK)
    add_textbox(slide, x + Inches(2.9), y + Pt(2), Inches(3.0), row_h, desc, size=10, color=PPTX_WARM_GRAY)

# ══════════════════════════════════════
#  3. 会议服务总则
# ══════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_section_title(slide, "01  会议服务总则", "GENERAL PRINCIPLES")
add_multiline(slide, Inches(0.6), Inches(1.8), Inches(7.5), Inches(5.0),
    ["一、适用范围",
     "本SOP适用于酒店内所有会议、宴会、商务活动及社交活动的服务操作。",
     "",
     "二、服务宗旨",
     "以专业化、标准化、个性化的服务，确保每一场会议圆满成功。",
     "",
     "三、职责划分",
     "• 会议销售部：客户对接、方案报价、合同签订",
     "• 宴会服务部：场地布置、现场服务、餐饮保障",
     "• 工程部：设备调试、电力保障、空调/灯光控制",
     "• 厨房部：茶歇/宴会出品、食品安全",
     "• 安保部：消防检查、车辆引导、秩序维护",
     "",
     "四、服务原则",
     "• 首问负责制：第一个接到需求的员工全程跟进",
     "• 15分钟响应制：客户需求15分钟内给出解决方案",
     "• 三级检查制：主管→经理→总监逐级复核",
    ],
    size=13, color=PPTX_DARK, spacing=1.3)

# 右下角卡片
add_card(slide, Inches(8.5), Inches(1.8), Inches(4.3), Inches(1.5),
         "服务热线", ["会议服务专线：400-XXX-XXXX",
                      "24小时值班经理：1XX-XXXX-XXXX"],
         title_size=16)

add_card(slide, Inches(8.5), Inches(3.6), Inches(4.3), Inches(1.5),
         "质量标准", ["ISO 9001 服务管理体系",
                      "金钥匙服务标准认证"],
         title_size=16)

# ══════════════════════════════════════
#  4. 会议预订流程
# ══════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_section_title(slide, "02  会议预订流程", "BOOKING PROCESS")

steps = [
    ("1", "客户咨询", "了解会议基本信息：\n规模/时间/预算/需求"),
    ("2", "方案制定", "出具会议方案书：\n场地/餐饮/设备报价"),
    ("3", "合同确认", "签订会议合同：\n支付定金/锁定档期"),
    ("4", "信息录入", "录入PMS/宴会系统：\n生成会议接待单"),
    ("5", "会前协调", "召开会前协调会：\n各部门确认分工"),
    ("6", "现场执行", "按SOP执行布场：\n管家全程跟进"),
    ("7", "结算归档", "费用结算/发票开具：\n客户档案归档"),
]
for i, (num, title, desc) in enumerate(steps):
    x = Inches(0.5 + i * 1.85)
    y = Inches(2.2)
    # arrow
    if i < len(steps) - 1:
        arr = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                      x + Inches(1.3), y + Inches(0.55),
                                      Inches(0.7), Pt(2))
        arr.fill.solid(); arr.fill.fore_color.rgb = PPTX_GOLD; arr.line.fill.background()
    shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, x + Inches(0.35), y, Inches(0.6), Inches(0.6))
    shape.fill.solid(); shape.fill.fore_color.rgb = PPTX_NAVY
    shape.line.color.rgb = PPTX_GOLD; shape.line.width = Pt(2)
    add_textbox(slide, x + Inches(0.35), y + Pt(5), Inches(0.6), Inches(0.5),
                num, size=18, bold=True, color=PPTX_GOLD, align=PP_ALIGN.CENTER)
    add_textbox(slide, x, y + Inches(0.7), Inches(1.7), Inches(0.35),
                title, size=13, bold=True, color=PPTX_DARK, align=PP_ALIGN.CENTER)
    add_textbox(slide, x - Inches(0.1), y + Inches(1.05), Inches(1.9), Inches(1.0),
                desc, size=10, color=PPTX_WARM_GRAY, align=PP_ALIGN.CENTER)

# ══════════════════════════════════════
#  5. 会前准备SOP
# ══════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_section_title(slide, "03  会前准备SOP", "PRE-EVENT PREPARATION")
add_multiline(slide, Inches(0.6), Inches(1.8), Inches(6.0), Inches(5.0),
    ["【T-3天】前期准备",
     "□ 接收会议接待通知单",
     "□ 确认台型/人数/时间/特殊要求",
     "□ 预留场地并锁定时段",
     "",
     "【T-1天】最终确认",
     "□ 与销售/客户确认最终人数",
     "□ 确认菜单/酒水/茶歇方案",
     "□ 检查设备运行状态（投影/音响/LED）",
     "□ 准备会议文具/水牌/席卡/签到台",
     "",
     "【T-3小时】现场布置",
     "□ 按标准台型图纸摆台",
     "□ 设备开机测试",
     "□ 空调提前开启（24-26°C）",
     "□ 灯光调试（会议模式/投影模式）",
     "",
     "【T-0.5小时】最终检查",
     "□ 主管逐项复核",
     "□ 会议管家到岗就位"],
    size=12, color=PPTX_DARK, spacing=1.2)

add_card(slide, Inches(7.0), Inches(1.8), Inches(5.8), Inches(2.0),
         "会前检查确认表", ["场地清洁度□  桌椅稳固性□  设备运行□",
                          "空调温度□   灯光模式□  音响测试□",
                          "水牌席卡□   会议文具□  签到设施□",
                          "茶歇准备□   卫生间清洁□ 安全出口指示□"],
         title_size=16, body_size=12)
add_card(slide, Inches(7.0), Inches(4.2), Inches(5.8), Inches(2.5),
         "注意事项", ["• 所有设备提前2小时完成调试",
                      "• 预留10%备用椅子/餐具/文具",
                      "• 同声传译设备需提前1天调试",
                      "• 会议临时需求30分钟内响应",
                      "• 禁止在会议区域使用手机"],
         title_size=16, body_size=12)

# ══════════════════════════════════════
#  6. 剧院式  (带图)
# ══════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_section_title(slide, "04  剧院式布局", "THEATER STYLE")
add_textbox(slide, Inches(0.6), Inches(1.6), Inches(5.5), Inches(0.3),
            "适用于：大型会议、演讲、培训、产品发布会", size=13, color=PPTX_WARM_GRAY)
slide.shapes.add_picture(layout_images["theater"],
                          Inches(6.8), Inches(1.0), width=Inches(6.0))
add_multiline(slide, Inches(0.6), Inches(2.1), Inches(5.8), Inches(4.5),
    ["【布局标准】",
    "• 椅子排列：每排间距 ≥ 45cm，排间距 ≥ 80cm",
    "• 走道设置：每列不超过14个座位，中间留纵走道",
    "• 讲台高度：30-40cm，配演讲台/立式麦克",
    "• 第一排距讲台 ≥ 2m",
    "",
    "【适合场景】",
    "• 大型主题演讲、年度总结大会",
    "• 新闻发布会、产品发布会",
    "• 公开课/培训讲座（纯听讲形式）",
    "• 颁奖典礼、开幕仪式",
    "",
    "【容纳参考】",
    "• 标准宴会厅 800-1200人",
    "• 中型厅 200-400人",
    "",
    "【注意事项】",
    "• 靠墙留边距 ≥ 60cm，方便通行",
    "• 后排可适当加高（阶梯式）"],
    size=11, color=PPTX_DARK, spacing=1.1)
add_page_num(slide, 6)

# ══════════════════════════════════════
#  7. 课桌式 (带图)
# ══════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_section_title(slide, "05  课桌式布局", "CLASSROOM STYLE")
add_textbox(slide, Inches(0.6), Inches(1.6), Inches(5.5), Inches(0.3),
            "适用于：培训、考试、研讨会、工作坊", size=13, color=PPTX_WARM_GRAY)
slide.shapes.add_picture(layout_images["classroom"],
                          Inches(6.8), Inches(1.0), width=Inches(6.0))
add_multiline(slide, Inches(0.6), Inches(2.1), Inches(5.8), Inches(4.5),
    ["【布局标准】",
    "• 课桌尺寸：180cm×45cm（标准），配1椅/位",
    "• 桌间距：排距 ≥ 90cm（方便进出）",
    "• 每排桌数：6-8张，中间留纵走道",
    "• 桌面摆水/纸笔/议程单（间距均匀）",
    "",
    "【适合场景】",
    "• 职业技能培训、认证考试",
    "• 研讨会、工作坊（需记笔记）",
    "• 小组讨论配合讲师讲授",
    "• 董事会报告（需翻阅资料）",
    "",
    "【容纳参考】",
    "• 每桌2人（一侧）或4人（两侧）",
    "• 500m²约容纳 150-200人",
    "",
    "【注意事项】",
    "• 电源排插提前布线（笔记本充电）",
    "• Wi-Fi密码打印在席卡上"],
    size=11, color=PPTX_DARK, spacing=1.1)
add_page_num(slide, 7)

# ══════════════════════════════════════
#  8. U型 (带图)
# ══════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_section_title(slide, "06  U型布局", "U-SHAPE STYLE")
add_textbox(slide, Inches(0.6), Inches(1.6), Inches(5.5), Inches(0.3),
            "适用于：商务洽谈、座谈会、联席会", size=13, color=PPTX_WARM_GRAY)
slide.shapes.add_picture(layout_images["ushape"],
                          Inches(6.8), Inches(1.0), width=Inches(6.0))
add_multiline(slide, Inches(0.6), Inches(2.1), Inches(5.8), Inches(4.5),
    ["【布局标准】",
    "• 桌子拼成长U形，开口朝向讲台/投影",
    "• 外侧坐人，内侧空出",
    "• 转角处用45°拼接桌或圆角桌",
    "• 桌面铺墨绿色绒布或白色桌布",
    "",
    "【适合场景】",
    "• 跨部门联席会议、协调会",
    "• 商务谈判、签约仪式",
    "• 学术委员会、评审会",
    "• 政府接待座谈会",
    "",
    "【容纳参考】",
    "• U型单侧可坐 8-15人",
    "• 中型厅约容纳 30-60人",
    "",
    "【注意事项】",
    "• 开口方向避开空调出风口",
    "• 桌面摆话筒（主席位+每位嘉宾）"],
    size=11, color=PPTX_DARK, spacing=1.1)
add_page_num(slide, 8)

# ══════════════════════════════════════
#  9. 回字型 (带图)
# ══════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_section_title(slide, "07  回字型布局", "HOLLOW SQUARE STYLE")
add_textbox(slide, Inches(0.6), Inches(1.6), Inches(5.5), Inches(0.3),
            "适用于：圆桌会议、多方会谈、董事会", size=13, color=PPTX_WARM_GRAY)
slide.shapes.add_picture(layout_images["hollow_square"],
                          Inches(6.8), Inches(1.0), width=Inches(6.0))
add_multiline(slide, Inches(0.6), Inches(2.1), Inches(5.8), Inches(4.5),
    ["【布局标准】",
    "• 桌子拼合为长方形或正方形闭环",
    "• 中间留空（可摆放绿植或会议主题装饰）",
    "• 内侧不坐人，外侧围坐",
    "• 四角用45°拼接桌过渡",
    "",
    "【适合场景】",
    "• 董事会会议、股东会议",
    "• 国际多方视频会议",
    "• 圆桌论坛（中间放话筒）",
    "• 高层战略研讨会",
    "",
    "【容纳参考】",
    "• 每边6-12人，共 20-50人",
    "• 建议最大不超过 60人",
    "",
    "【注意事项】",
    "• 中间摆放鲜花或冰雕（高度不超过视线）",
    "• 同声传译耳机接口每座配备"],
    size=11, color=PPTX_DARK, spacing=1.1)
add_page_num(slide, 9)

# ══════════════════════════════════════
#  10. 鱼骨式 (带图)
# ══════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_section_title(slide, "08  鱼骨式布局", "FISHBONE STYLE")
add_textbox(slide, Inches(0.6), Inches(1.6), Inches(5.5), Inches(0.3),
            "适用于：互动培训、头脑风暴、团队共创", size=13, color=PPTX_WARM_GRAY)
slide.shapes.add_picture(layout_images["fishbone"],
                          Inches(6.8), Inches(1.0), width=Inches(6.0))
add_multiline(slide, Inches(0.6), Inches(2.1), Inches(5.8), Inches(4.5),
    ["【布局标准】",
    "• 矩形桌以中轴为基准，左右45°斜摆",
    "• 形似鱼骨，对向讲台/主屏幕",
    "• 每桌坐4-6人（单侧坐人）",
    "• 中轴线留宽走道（≥1.5m）",
    "",
    "【适合场景】",
    "• 互动型培训、团队工作坊",
    "• 头脑风暴、创意策划会",
    "• 分组讨论+集中汇报模式",
    "• 敏捷开发冲刺会议",
    "",
    "【容纳参考】",
    "• 每侧4-6组，共 40-80人",
    "",
    "【注意事项】",
    "• 确保每桌学员能看到主屏幕",
    "• 走道铺地毯降噪"],
    size=11, color=PPTX_DARK, spacing=1.1)
add_page_num(slide, 10)

# ══════════════════════════════════════
#  11. 宴会式 (带图)
# ══════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_section_title(slide, "09  宴会式布局", "BANQUET STYLE")
add_textbox(slide, Inches(0.6), Inches(1.6), Inches(5.5), Inches(0.3),
            "适用于：宴会、年会、颁奖晚宴、社交活动", size=13, color=PPTX_WARM_GRAY)
slide.shapes.add_picture(layout_images["banquet"],
                          Inches(6.8), Inches(1.0), width=Inches(6.0))
add_multiline(slide, Inches(0.6), Inches(2.1), Inches(5.8), Inches(4.5),
    ["【布局标准】",
    "• 圆桌直径：1.8m（10人）/ 2.2m（12人）",
    "• 桌间距 ≥ 2m（方便服务员上菜）",
    "• 主桌位于舞台正前方（直径2.8m）",
    "• 每桌配转盘、口布折花、菜单卡",
    "",
    "【适合场景】",
    "• 公司年会、答谢宴、团圆宴",
    "• 颁奖典礼+晚宴",
    "• 婚宴、寿宴等社交宴会",
    "• 慈善拍卖晚宴",
    "",
    "【容纳参考】",
    "• 800m²约可设 25-30桌（10人/桌）",
    "• 大型宴会厅可达 80-100桌",
    "",
    "【注意事项】",
    "• 预留舞台前方摄影通道",
    "• 主桌设专人服务（分菜/斟酒）"],
    size=11, color=PPTX_DARK, spacing=1.1)
add_page_num(slide, 11)

# ══════════════════════════════════════
#  12. 鸡尾酒会式 (带图)
# ══════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_section_title(slide, "10  鸡尾酒会式", "COCKTAIL STYLE")
add_textbox(slide, Inches(0.6), Inches(1.6), Inches(5.5), Inches(0.3),
            "适用于：交流酒会、开幕派对、商务社交", size=13, color=PPTX_WARM_GRAY)
slide.shapes.add_picture(layout_images["cocktail"],
                          Inches(6.8), Inches(1.0), width=Inches(6.0))
add_multiline(slide, Inches(0.6), Inches(2.1), Inches(5.8), Inches(4.5),
    ["【布局标准】",
    "• 高脚圆桌（高度1.1m）分散布置",
    "• 每桌配4-6张高脚椅或站立式",
    "• 沿墙设吧台/取餐台/酒水台",
    "• 场地中央保持通畅",
    "",
    "【适合场景】",
    "• 开幕酒会、媒体见面会",
    "• 展览开幕式、画廊开幕",
    "• 商务社交、行业交流活动",
    "• 签约仪式后的庆祝酒会",
    "",
    "【容纳参考】",
    "• 按 1.5-2m²/人 计算",
    "• 灵活度高，可达数百人",
    "",
    "【注意事项】",
    "• 地面防滑处理（酒水易洒落）",
    "• 服务员托盘穿梭服务"],
    size=11, color=PPTX_DARK, spacing=1.1)
add_page_num(slide, 12)

# ══════════════════════════════════════
#  13. 董事会式 (带图)
# ══════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_section_title(slide, "11  董事会式布局", "BOARDROOM STYLE")
add_textbox(slide, Inches(0.6), Inches(1.6), Inches(5.5), Inches(0.3),
            "适用于：董事会议、高层决策会、VIP洽谈", size=13, color=PPTX_WARM_GRAY)
slide.shapes.add_picture(layout_images["boardroom"],
                          Inches(6.8), Inches(1.0), width=Inches(6.0))
add_multiline(slide, Inches(0.6), Inches(2.1), Inches(5.8), Inches(4.5),
    ["【布局标准】",
    "• 大型椭圆/长方形会议桌",
    "• 真皮高背椅，每座间隔 ≥ 80cm",
    "• 每位配液晶显示屏/接口面板",
    "• 桌面摆：皮垫/水杯/笔记本/签字笔",
    "",
    "【适合场景】",
    "• 上市公司董事会会议",
    "• 跨国集团高层战略会议",
    "• VIP客户闭门洽谈",
    "• 投资决策委员会会议",
    "",
    "【容纳参考】",
    "• 标准董事会室 14-30人",
    "• 建议不超过 40人",
    "",
    "【注意事项】",
    "• 提供同声传译（国际会议）",
    "• 会议保密协议签署"],
    size=11, color=PPTX_DARK, spacing=1.1)
add_page_num(slide, 13)

# ══════════════════════════════════════
#  14. T型台 (带图)
# ══════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_section_title(slide, "12  T型台布局", "T-SHAPE STYLE")
add_textbox(slide, Inches(0.6), Inches(1.6), Inches(5.5), Inches(0.3),
            "适用于：时装秀、产品展示、舞台表演", size=13, color=PPTX_WARM_GRAY)
slide.shapes.add_picture(layout_images["tshape"],
                          Inches(6.8), Inches(1.0), width=Inches(6.0))
add_multiline(slide, Inches(0.6), Inches(2.1), Inches(5.8), Inches(4.5),
    ["【布局标准】",
    "• T台宽度：1.2-2.4m（视活动类型）",
    "• T台高度：30-60cm",
    "• 两侧剧院式座椅",
    "• 台面铺黑色地毯或镜面地板",
    "",
    "【适合场景】",
    "• 时装秀、新品发布走秀",
    "• 产品实物展示演示",
    "• 大型舞台表演活动",
    "• 颁奖仪式（走红毯环节）",
    "",
    "【容纳参考】",
    "• 两侧各6-8排，共 100-300人",
    "",
    "【注意事项】",
    "• T台承重检查（走秀/重物）",
    "• 灯光追光系统配合"],
    size=11, color=PPTX_DARK, spacing=1.1)
add_page_num(slide, 14)

# ══════════════════════════════════════
#  15. 台型选择指南
# ══════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_section_title(slide, "13  台型选择指南", "LAYOUT SELECTION GUIDE")

headers = ["台型", "容纳能力", "互动性", "适合活动类型", "建议桌型"]
rows = [
    ["剧院式", "80-1200人", "低", "演讲/发布会/年会", "无桌/椅子"],
    ["课桌式", "50-400人", "中", "培训/考试/研讨会", "长条桌"],
    ["U型", "20-60人", "高", "洽谈/座谈会/协调会", "长条桌拼接"],
    ["回字型", "20-50人", "高", "董事会/多方会谈", "长条桌+角桌"],
    ["鱼骨式", "30-80人", "很高", "工作坊/头脑风暴", "长条桌斜摆"],
    ["宴会式", "50-1000人", "中低", "晚宴/年会/婚宴", "圆桌"],
    ["鸡尾酒会", "50-500人", "高", "酒会/社交/开幕", "高脚圆桌"],
    ["董事会式", "10-30人", "很高", "高层决策/董事会议", "椭圆桌"],
    ["T型台", "100-300人", "低", "时装秀/走秀/展示", "T台+座椅"],
]
col_widths = [Inches(1.4), Inches(1.6), Inches(1.2), Inches(3.8), Inches(2.2)]
add_sop_table(slide, Inches(0.6), Inches(1.8), col_widths, headers, rows)

# 补充说明
add_card(slide, Inches(0.6), Inches(5.2), Inches(5.8), Inches(1.8),
         "选型建议", ["• 听讲为主 → 剧院式（性价比最高）",
                      "• 需记录 → 课桌式、鱼骨式",
                      "• 互动讨论 → U型、回字型",
                      "• 正式餐宴 → 宴会式（圆桌）"],
         title_size=15, body_size=11)
add_card(slide, Inches(6.8), Inches(5.2), Inches(5.8), Inches(1.8),
         "台型转换时间参考", ["• 剧院式→课桌式：约45分钟（4人）",
                              "• 课桌式→宴会式：约90分钟（6人）",
                              "• U型→回字型：约30分钟（2人）",
                              "• 全清重新摆台：约120分钟（8人）"],
         title_size=15, body_size=11)

# ══════════════════════════════════════
#  16. 会议设备标准 (带图)
# ══════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_section_title(slide, "14  会议设备标准", "AV EQUIPMENT STANDARDS")
slide.shapes.add_picture(layout_images["equipment"],
                          Inches(0.5), Inches(1.6), width=Inches(6.5))

headers2 = ["设备", "规格要求", "检查频次", "责任人"]
rows2 = [
    ["投影仪", "≥5000流明，支持1080P", "每次使用前", "工程部"],
    ["投影幕布", "16:9 电动幕，≥150寸", "每周清洁一次", "工程部"],
    ["LED大屏", "P2.5/P3全彩，含处理器", "每月校色一次", "工程部"],
    ["音响系统", "线阵/柱式，覆盖全场", "每次使用前测试", "工程部"],
    ["无线麦克", "≥4支手持+2支领夹", "每次检查电池", "宴会部"],
    ["同声传译", "红外/射频，≥4语种", "每次活动前调试", "外租/工程"],
    ["视频会议", "4K摄像头+全向麦", "每次连线前测试", "工程部"],
    ["灯光系统", "LED染色灯/追光/面光", "每月维护一次", "工程部"],
]
col_widths2 = [Inches(1.8), Inches(3.0), Inches(1.8), Inches(1.6)]
add_sop_table(slide, Inches(7.2), Inches(1.6), col_widths2, headers2, rows2)

add_card(slide, Inches(7.2), Inches(5.2), Inches(5.5), Inches(1.8),
         "设备故障应急方案", ["• 投影故障 → 备用投影仪5分钟到位",
                              "• 麦克无声 → 立即换电池/换备用麦",
                              "• 音响啸叫 → 工程部30秒内调参",
                              "• LED黑屏 → 备用信源切换+工程抢修"],
         title_size=15, body_size=11)

# ══════════════════════════════════════
#  17. 会议摆台标准
# ══════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_section_title(slide, "15  会议摆台标准", "TABLE SETUP STANDARDS")

# 左侧：座位摆台
add_card(slide, Inches(0.5), Inches(1.8), Inches(5.8), Inches(2.8),
         "会议座位摆台标准（每位）", [
             "① 会议垫/文件夹 — 居中，距桌边1cm",
             "② 签字笔（黑色）— 放在文件夹右侧",
             "③ 会议议程单 — 文件夹内左侧",
             "④ 会议记录纸（5张）— 文件夹内右侧",
             "⑤ 矿泉水 — 文件夹左前方45°",
             "⑥ 玻璃杯 — 矿泉水左侧",
             "⑦ 纸杯+杯垫 — 文件夹右前方45°",
             "⑧ 席卡（如有）— 文件夹正前方",
         ],
         title_size=14, body_size=11)

add_card(slide, Inches(0.5), Inches(4.9), Inches(5.8), Inches(2.0),
         "布场质量验收标准", [
             "• 椅子对齐：横竖成直线（用拉线校验）",
             "• 桌布垂边：统一25cm，无褶皱",
             "• 桌面物品：Logo朝向统一（正对客人）",
             "• 走道宽度：主通道≥1.5m，次通道≥1.0m",
         ],
         title_size=14, body_size=11)

# 右侧：水牌/席卡
add_card(slide, Inches(6.8), Inches(1.8), Inches(5.8), Inches(2.0),
         "会议水牌标准", [
             "• 尺寸：300mm×400mm（立式）",
             "• 内容：会议名称 + 主办方 + 日期",
             "• 材质：亚克力/金属/覆膜纸质",
             "• 摆放位置：会场入口两侧 + 电梯口",
             "• 英文在上，中文在下（国际会议）",
         ],
         title_size=14, body_size=11)

add_card(slide, Inches(6.8), Inches(4.1), Inches(5.8), Inches(2.8),
         "席卡（台卡）标准", [
             "• 尺寸：100mm×200mm（双面显示）",
             "• 字体：黑体/微软雅黑，字号≥100pt",
             "• 内容：客人姓名 + 职务（中英文）",
             "• 材质：亚克力/透明插卡/铜版纸折卡",
             "• 摆放：正对座位方向，居中放置",
             "",
             "【席卡排列原则】",
             "• 主席位：面朝入口/屏幕的中央位置",
             "• 客方最高领导坐主人右手边",
             "• 中文姓名：姓在前，名在后（标准格式）",
         ],
         title_size=14, body_size=11)

# ══════════════════════════════════════
#  18. 茶歇服务SOP (带图)
# ══════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_section_title(slide, "16  茶歇服务SOP", "TEA BREAK SERVICE")
slide.shapes.add_picture(layout_images["tea_setup"],
                          Inches(6.5), Inches(1.2), width=Inches(6.5))
add_multiline(slide, Inches(0.5), Inches(1.8), Inches(5.8), Inches(5.0),
    ["【摆台标准】",
    "• 咖啡/茶饮区：咖啡机2台+茶包8种+热水",
    "• 甜品区：甜点3种+水果拼盘2种+果切",
    "• 冷餐区：三明治/迷你汉堡/沙拉杯",
    "• 餐具区：杯碟+纸巾+搅拌棒+糖包奶盅",
    "",
    "【服务流程】",
    "① 茶歇前30分钟：完成摆台，检查食品温度",
    "② 茶歇前15分钟：补充饮品（咖啡/茶现做）",
    "③ 会议休息：服务员站位引导（微笑问候）",
    "④ 中进行中：巡视补充（食品低于50%即补）",
    "⑤ 结束后：15分钟内撤台清洁",
    "",
    "【卫生标准】",
    "• 所有食品加盖/覆保鲜膜直至服务前",
    "• 饮品容器每2小时更换清洗",
    "• 服务员戴手套+口罩操作",
    "• 过敏源信息标注（含坚果/乳制品等）",
    "",
    "【时段标准】",
    "• 上午茶歇：10:00-10:15（15分钟）",
    "• 下午茶歇：15:00-15:15（15分钟）",
    "• 全天会议：上下午各一次 + 午餐自助"],
    size=11, color=PPTX_DARK, spacing=1.1)

# ══════════════════════════════════════
#  19. 会议餐饮服务
# ══════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_section_title(slide, "17  会议餐饮服务", "CATERING SERVICE")

add_card(slide, Inches(0.5), Inches(1.8), Inches(3.8), Inches(2.5),
         "自助午餐标准", [
             "• 热菜：4荤+4素+2主食",
             "• 沙拉吧：6种蔬菜+4种酱汁",
             "• 汤品：2种（中西各一）",
             "• 水果：3种时令水果",
             "• 饮品：果汁+软饮+咖啡茶",
             "• 服务：80人以内1位厨师巡场",
             "• 布菲台：每80人设一台",
         ],
         title_size=14, body_size=10)

add_card(slide, Inches(4.6), Inches(1.8), Inches(3.8), Inches(2.5),
         "中式宴会标准", [
             "• 冷菜：6-8道精美冷碟",
             "• 热菜：8-10道（含海鲜/汤/甜品）",
             "• 主食：炒饭/面点/点心组合",
             "• 酒水：白酒/红酒/软饮/茶",
             "• 上菜节奏：凉菜→热菜→主食→甜品",
             "• 每道菜间隔：5-8分钟",
             "• 主桌设分菜服务",
         ],
         title_size=14, body_size=10)

add_card(slide, Inches(8.7), Inches(1.8), Inches(4.0), Inches(2.5),
         "西式/酒会标准", [
             "• Canapé小食：8-10种（热/冷各半）",
             "• 芝士拼盘：4种芝士+干果+蜂蜜",
             "• 现做档口：煎鹅肝/生蚝/牛排",
             "• 酒水：香槟/红酒/白葡萄酒/鸡尾酒",
             "• 服务方式：托盘穿梭+驻站吧台",
             "• 服务员配比：1:30客人",
         ],
         title_size=14, body_size=10)

add_multiline(slide, Inches(0.5), Inches(4.6), Inches(12.3), Inches(2.5),
    ["【餐饮服务通用标准】",
    "• 所有菜品测试：活动前3天进行试菜（主家确认）",
    "• 过敏源管理：菜单标注含麸质/坚果/乳制品/海鲜",
    "• 备餐比例：正餐备 110%（含服务人员）",
    "• 上菜时间误差：不超过 ±5分钟",
    "• 清真餐/素食/特殊餐：提前48小时沟通确认",
    "• 食品留样：每道菜125g，冷藏保存48小时"],
    size=12, color=PPTX_DARK, spacing=1.3)

# ══════════════════════════════════════
#  20. 会议管家职责
# ══════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_section_title(slide, "18  会议管家职责", "EVENT MANAGER DUTIES")

add_card(slide, Inches(0.5), Inches(1.8), Inches(3.8), Inches(2.0),
         "岗前准备", [
             "• 提前3天对接销售/客户",
             "• 建立会议专属服务群",
             "• 确认会议流程单",
             "• 协调各部门分工",
         ],
         title_size=14, body_size=12)

add_card(slide, Inches(4.6), Inches(1.8), Inches(3.8), Inches(2.0),
         "会中执行", [
             "• 全程在会场值守",
             "• 每30分钟巡视一次",
             "• 处理客户临时需求",
             "• 协调设备/餐饮/清洁",
         ],
         title_size=14, body_size=12)

add_card(slide, Inches(8.7), Inches(1.8), Inches(4.0), Inches(2.0),
         "会后跟进", [
             "• 协助结账/开票",
             "• 收集客户反馈",
             "• 整理会议档案",
             "• 建立客户服务档案",
         ],
         title_size=14, body_size=12)

add_multiline(slide, Inches(0.5), Inches(4.2), Inches(12.3), Inches(2.8),
    ["【会议管家服务动线】",
    "07:30 到岗 → 检查会场/设备/茶歇",
    "08:00 迎宾（门口站位）→ 引导签到/指引入座",
    "08:30 会议开始 → 关注设备运行/茶水服务",
    "10:00 茶歇服务 → 引导客人至茶歇区",
    "10:15 会议继续 → 快速整理会场/补充文具",
    "12:00 午餐 → 引导至餐厅/安排分菜",
    "13:30 下午会议 → 同上午流程",
    "15:00 下午茶歇 → 同上",
    "17:00 会议结束 → 协助送客/收集反馈/场地复原",
    "",
    "【注意事项】",
    "• 会议管家手机静音，仅振动模式",
    "• 站立姿势：双手交叠自然下垂，不靠墙",
    "• 与客户交流：先微笑，再说话，语速适中",
    "• 客人提问不推诿：'我来为您解决' 为首要回应"],
    size=11, color=PPTX_DARK, spacing=1.2)

# ══════════════════════════════════════
#  21. 应急预案
# ══════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_section_title(slide, "19  应急预案", "EMERGENCY RESPONSE")

headers3 = ["突发情况", "应急措施", "责任人"]
rows3 = [
    ["设备故障", "立即启用备用设备，5分钟内响应", "工程部值班员"],
    ["停电/跳闸", "启动UPS应急照明，30秒内切换备用电源", "工程部主管"],
    ["麦克啸叫", "立即降低音量+调整EQ，30秒内解决", "音响师"],
    ["客人突发疾病", "立即拨打120+酒店医务室+报告值班经理", "会议管家"],
    ["食品过敏", "立即停止食用+呼叫医护+保留样本", "宴会部+厨房"],
    ["火灾警报", "按消防预案引导疏散（会议管家领队）", "安保部"],
    ["客人投诉", "立即道歉+记录+报告经理+15分钟方案", "会议管家"],
    ["临时增加人数", "30分钟内准备备用桌椅/餐具/餐食", "宴会部"],
]
col_widths3 = [Inches(2.5), Inches(5.5), Inches(2.2)]
add_sop_table(slide, Inches(0.5), Inches(1.8), col_widths3, headers3, rows3)

add_card(slide, Inches(0.5), Inches(5.2), Inches(12.3), Inches(1.8),
         "应急物资储备", [
             "• 备用投影仪×1  • 延长线排插×10  • 备用麦克风电池（5号/9V）×各20",
             "• 急救箱（含常用药/创可贴/过敏药）×2  • 备用席卡/水牌空白模板×30",
             "• 应急照明手电×10  • 对讲机（每名管家/安保/工程各配1台）",
         ],
         title_size=14, body_size=12)

# ══════════════════════════════════════
#  22. 会议物品清单
# ══════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_section_title(slide, "20  会议物品清单", "EQUIPMENT CHECKLIST")

headers4 = ["类别", "物品名称", "规格/数量", "备注"]
rows4 = [
    ["摆台工具", "钢卷尺/拉线绳", "各2", "校验桌椅对齐"],
    ["摆台工具", "水平仪", "1把", "检查桌面水平"],
    ["文具", "会议笔（黑）", "按人数×2", "备用"],
    ["文具", "会议记录纸/A4纸", "按人数×5张", ""],
    ["文具", "文件夹/会议垫", "按人数", ""],
    ["文具", "白板笔（黑/蓝/红）", "各5支", "含板擦"],
    ["标识", "席卡/台卡", "按人数+20%备用", "打印模板"],
    ["标识", "会议水牌", "2-4个", "入口/电梯"],
    ["标识", "签到板/背景板", "1套", "按设计制作"],
    ["饮品", "矿泉水", "按人数×4瓶", "会议全天"],
    ["饮品", "咖啡/茶包", "按人数×2份", "茶歇用"],
    ["服务", "托盘", "6个", "茶水服务"],
    ["服务", "对讲机", "按岗位", "统一频道"],
    ["应急", "备用电池", "5号20节/9V10节", "麦克用"],
    ["应急", "急救箱", "1个", "前台领用"],
]
col_widths4 = [Inches(1.2), Inches(2.5), Inches(2.0), Inches(5.5)]
add_sop_table(slide, Inches(0.5), Inches(1.8), col_widths4, headers4, rows4)

# ══════════════════════════════════════
#  23. 附录
# ══════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_section_title(slide, "21  附录", "APPENDIX")

add_card(slide, Inches(0.5), Inches(1.8), Inches(5.8), Inches(2.5),
         "附件一：会议接待通知单", [
             "• 会议名称/主办方/联系人",
             "• 日期/时间/场地/人数/台型",
             "• 设备/餐饮/住宿要求",
             "• 特殊要求/VIP接待方案",
             "• 费用预算/结算方式",
         ],
         title_size=14, body_size=12)

add_card(slide, Inches(6.8), Inches(1.8), Inches(5.8), Inches(2.5),
         "附件二：会前检查表", [
             "□ 场地清洁 □ 台型确认 □ 设备测试",
             "□ 文具到位 □ 席卡打印 □ 水牌摆放",
             "□ 签到台设置 □ 茶歇准备 □ 空调开启",
             "□ 灯光调试 □ 音响测试 □ 投影检查",
             "□ 安全出口指示 □ 卫生间清洁 □ 绿植装饰",
         ],
         title_size=14, body_size=12)

add_card(slide, Inches(0.5), Inches(4.6), Inches(5.8), Inches(2.5),
         "附件三：会议服务质量评价表", [
             "• 场地布置满意度（1-5分）",
             "• 设备运行满意度（1-5分）",
             "• 餐饮质量满意度（1-5分）",
             "• 服务态度满意度（1-5分）",
             "• 综合评价与改进建议",
         ],
         title_size=14, body_size=12)

add_card(slide, Inches(6.8), Inches(4.6), Inches(5.8), Inches(2.5),
         "附件四：会议交接班记录表", [
             "• 会议名称/日期/当前进度",
             "• 已完成事项/待完成事项",
             "• 设备运行状态/异常记录",
             "• 客户特殊需求/注意事项",
             "• 交班人/接班人签名",
         ],
         title_size=14, body_size=12)

# ══════════════════════════════════════
#  24. 封底
# ══════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_shape_bg(slide, PPTX_NAVY)
add_gold_bar(slide, Inches(5.0), Inches(2.2), Inches(3.3), Pt(2))
add_textbox(slide, Inches(1), Inches(2.5), Inches(11.3), Inches(0.8),
            "标准成就卓越", size=40, bold=True, color=PPTX_GOLD, align=PP_ALIGN.CENTER)
add_textbox(slide, Inches(1), Inches(3.4), Inches(11.3), Inches(0.5),
            "细节决定品质  专业成就完美", size=20, color=PPTX_WHITE, align=PP_ALIGN.CENTER)
add_gold_bar(slide, Inches(5.0), Inches(4.0), Inches(3.3), Pt(2))
add_multiline(slide, Inches(2), Inches(4.5), Inches(9.3), Inches(1.5),
              ["[酒店名称]  会议服务部",
               "版本号：V1.0",
               "本手册为内部培训资料，未经授权不得外传"],
              size=14, color=PPTX_WARM_GRAY, align=PP_ALIGN.CENTER, spacing=1.5)

# ── Save ──
output_path = os.path.join(PPTX_DIR, "会议服务SOP手册.pptx")
prs.save(output_path)
print(f"[OK] SOP手册已生成: {output_path}")
print(f"[OK] 共 {len(prs.slides)} 页")
